from __future__ import annotations

import math
import shutil
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import (
    Image as RLImage,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)
from reportlab.platypus.tableofcontents import TableOfContents


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "report_output_rebuild"
ASSET_DIR = OUT / "assets"
FINAL_DIR = OUT / "final"
TMP_DIR = OUT / "tmp"
OLD_ASSET_DIR = ROOT / "docs" / "report_output" / "assets"

REPORT_DOCX = FINAL_DIR / "喵伴空气管家作品设计报告_重做版_2026-07-08.docx"
REPORT_PDF = FINAL_DIR / "喵伴空气管家作品设计报告_重做版_2026-07-08.pdf"
PPTX_PATH = FINAL_DIR / "喵伴空气管家答辩PPT_重做版_2026-07-08.pptx"
VIDEO_GUIDE = FINAL_DIR / "视频拍摄与答辩讲述建议_重做版_2026-07-08.md"

FONT_BODY = Path("C:/Windows/Fonts/STSONG.TTF")
FONT_BOLD = Path("C:/Windows/Fonts/simhei.ttf")
FONT_UI = Path("C:/Windows/Fonts/msyh.ttc")
FONT_UI_BOLD = Path("C:/Windows/Fonts/msyhbd.ttc")

GREEN = "#174238"
GREEN_2 = "#2E6B58"
LIGHT = "#F5F6F3"
ORANGE = "#D9822B"
GRAY = "#59615D"
DARK = "#1F2723"
LINE = "#C8D0C9"


@dataclass
class Figure:
    key: str
    title: str
    path: Path


@dataclass
class TableData:
    title: str
    headers: list[str]
    rows: list[list[str]]


@dataclass
class Section:
    title: str
    paragraphs: list[str]
    figure_key: str | None = None
    table: TableData | None = None


@dataclass
class Chapter:
    title: str
    sections: list[Section] = field(default_factory=list)


def ensure_dirs() -> None:
    for p in [ASSET_DIR, FINAL_DIR, TMP_DIR]:
        p.mkdir(parents=True, exist_ok=True)
    for name in [
        "hardware_overview_1.jpg",
        "hardware_overview_2.jpg",
        "serial_page_home.png",
        "serial_page_air_score.png",
        "serial_page_ai_settings.png",
        "upload_requirements.png",
    ]:
        src = OLD_ASSET_DIR / name
        dst = ASSET_DIR / name
        if src.exists() and not dst.exists():
            shutil.copyfile(src, dst)


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_UI_BOLD if bold and FONT_UI_BOLD.exists() else FONT_UI
    if not path.exists():
        path = FONT_BOLD if bold else FONT_BODY
    return ImageFont.truetype(str(path), size=size)


def text_size(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont) -> tuple[int, int]:
    box = draw.textbbox((0, 0), text, font=font)
    return box[2] - box[0], box[3] - box[1]


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    lines: list[str] = []
    current = ""
    for ch in text:
        if ch == "\n":
            if current:
                lines.append(current)
            current = ""
            continue
        candidate = current + ch
        if text_size(draw, candidate, font)[0] <= max_width or not current:
            current = candidate
        else:
            lines.append(current)
            current = ch
    if current:
        lines.append(current)
    return lines


def draw_wrapped(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    font: ImageFont.FreeTypeFont,
    fill: str,
    max_width: int,
    line_gap: int = 6,
    center: bool = False,
) -> int:
    x, y = xy
    for line in wrap_text(draw, text, font, max_width):
        w, h = text_size(draw, line, font)
        tx = x + (max_width - w) // 2 if center else x
        draw.text((tx, y), line, font=font, fill=fill)
        y += h + line_gap
    return y


def rounded_box(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    fill: str,
    outline: str = LINE,
    width: int = 2,
    radius: int = 18,
) -> None:
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], fill: str = GREEN, width: int = 4) -> None:
    draw.line([start, end], fill=fill, width=width)
    sx, sy = start
    ex, ey = end
    angle = math.atan2(ey - sy, ex - sx)
    head = 14
    pts = [
        (ex, ey),
        (ex - head * math.cos(angle - math.pi / 6), ey - head * math.sin(angle - math.pi / 6)),
        (ex - head * math.cos(angle + math.pi / 6), ey - head * math.sin(angle + math.pi / 6)),
    ]
    draw.polygon(pts, fill=fill)


def new_canvas(width: int = 1600, height: int = 900) -> tuple[Image.Image, ImageDraw.ImageDraw]:
    img = Image.new("RGB", (width, height), LIGHT)
    return img, ImageDraw.Draw(img)


def title(draw: ImageDraw.ImageDraw, text: str) -> None:
    draw.text((70, 46), text, font=load_font(44, True), fill=GREEN)
    draw.rectangle((70, 106, 230, 114), fill=ORANGE)


def save(img: Image.Image, name: str) -> Path:
    path = ASSET_DIR / name
    img.save(path, quality=95)
    return path


def draw_system_architecture() -> Path:
    img, draw = new_canvas()
    title(draw, "系统总体架构")
    layers = [
        ("感知与执行层", ["DHT11 温湿度", "MQ135 空气质量", "净化 LED", "加湿 LED", "舵机扇叶"]),
        ("ESP32S3 主控层", ["采样", "评分", "状态缓存", "自动/节能规则", "HTTP API"]),
        ("交互层", ["TJC 串口屏", "微信小程序", "小智 AI/MCP"]),
        ("应用层", ["家庭", "宿舍", "教室", "办公室"]),
    ]
    y = 150
    for i, (head, items) in enumerate(layers):
        x = 80 + i * 380
        rounded_box(draw, (x, y, x + 310, y + 575), "#FFFFFF", GREEN if i == 1 else LINE, 4 if i == 1 else 2, 24)
        draw.text((x + 28, y + 28), head, font=load_font(29, True), fill=GREEN)
        yy = y + 95
        for item in items:
            rounded_box(draw, (x + 30, yy, x + 280, yy + 62), "#F8FAF7", LINE, 2, 12)
            draw.text((x + 54, yy + 17), item, font=load_font(24), fill=DARK)
            yy += 84
        if i < len(layers) - 1:
            arrow(draw, (x + 318, y + 285), (x + 370, y + 285), ORANGE, 5)
    draw.text((90, 775), "主线：环境数据进入 ESP32S3，状态同步到屏幕、小程序和语音入口，控制结果再回到执行器。", font=load_font(27), fill=GRAY)
    return save(img, "01_system_architecture.png")


def draw_data_control_flow() -> Path:
    img, draw = new_canvas()
    title(draw, "数据流与控制流")
    center = (800, 450)
    draw.ellipse((610, 260, 990, 640), fill="#FFFFFF", outline=GREEN, width=6)
    draw_wrapped(draw, (675, 360), "ESP32S3\n状态计算与设备控制", load_font(34, True), GREEN, 250, center=True)
    nodes = [
        ("传感器采样", 190, 170),
        ("串口屏显示", 1170, 160),
        ("小程序读写", 1200, 570),
        ("AI/MCP 工具", 200, 575),
    ]
    for label, x, y in nodes:
        rounded_box(draw, (x, y, x + 245, y + 105), "#FFFFFF", LINE, 3, 18)
        draw_wrapped(draw, (x + 24, y + 28), label, load_font(28, True), DARK, 200, center=True)
    arrow(draw, (445, 222), (610, 360), ORANGE, 5)
    arrow(draw, (990, 360), (1170, 220), GREEN, 5)
    arrow(draw, (1170, 620), (990, 540), ORANGE, 5)
    arrow(draw, (445, 625), (610, 540), ORANGE, 5)
    labels = [
        ("温度/湿度/空气评分", 420, 250),
        ("控件刷新/曲线", 1015, 250),
        ("HTTP API", 1060, 640),
        ("语音转工具调用", 380, 640),
    ]
    for txt, x, y in labels:
        draw.text((x, y), txt, font=load_font(22), fill=GRAY)
    draw.text((96, 770), "控制流和数据流共用同一份设备状态，避免屏幕、手机、语音三个入口互相打架。", font=load_font(27), fill=GRAY)
    return save(img, "02_data_control_flow.png")


def draw_hardware_wiring() -> Path:
    img, draw = new_canvas()
    title(draw, "硬件连接示意")
    rounded_box(draw, (610, 220, 990, 620), "#FFFFFF", GREEN, 5, 32)
    draw.text((700, 285), "ESP32S3", font=load_font(46, True), fill=GREEN)
    pins = [
        ("GPIO18", "DHT11 温湿度 DATA", 115, 205, 610, 310),
        ("GPIO1 / ADC1_CH0", "MQ135 模拟输出", 115, 360, 610, 405),
        ("GPIO41/42", "TJC 串口屏 UART2", 1090, 205, 990, 310),
        ("GPIO13", "红色 LED 净化", 1090, 360, 990, 405),
        ("GPIO14", "蓝色 LED 加湿", 1090, 515, 990, 500),
        ("GPIO21", "舵机扇叶 新风", 115, 515, 610, 500),
    ]
    for pin, desc, x, y, sx, sy in pins:
        rounded_box(draw, (x, y, x + 380, y + 90), "#FFFFFF", LINE, 3, 16)
        draw.text((x + 24, y + 17), pin, font=load_font(25, True), fill=ORANGE)
        draw.text((x + 24, y + 50), desc, font=load_font(22), fill=DARK)
        ex = x + 380 if x < 600 else x
        arrow(draw, (sx, sy), (ex, y + 45), GREEN if x > 600 else ORANGE, 4)
    draw.text((105, 750), "供电和共地在实物接线中要单独确认，舵机建议使用稳定 5V 供电并与 ESP32S3 共地。", font=load_font(27), fill=GRAY)
    return save(img, "03_hardware_wiring.png")


def draw_serial_screen_pages() -> Path:
    img, draw = new_canvas()
    title(draw, "串口屏页面结构")
    pages = [
        ("首页", "温湿度、空气状态、建议"),
        ("空气详情", "空气评分、原始值、曲线"),
        ("智能家居", "净化、新风、加湿、自动、节能"),
        ("AI 与设置", "AI 状态、手动环境场景"),
    ]
    coords = [(130, 210), (940, 210), (130, 545), (940, 545)]
    for (name, desc), (x, y) in zip(pages, coords):
        rounded_box(draw, (x, y, x + 530, y + 210), "#FFFFFF", GREEN, 4, 26)
        draw.rectangle((x, y, x + 530, y + 60), fill=GREEN)
        draw.text((x + 28, y + 14), name, font=load_font(30, True), fill="#FFFFFF")
        draw_wrapped(draw, (x + 32, y + 90), desc, load_font(26), DARK, 455)
    arrow(draw, (660, 310), (940, 310), ORANGE, 5)
    arrow(draw, (940, 650), (660, 650), ORANGE, 5)
    arrow(draw, (395, 420), (395, 545), GREEN, 5)
    arrow(draw, (1205, 420), (1205, 545), GREEN, 5)
    draw.text((110, 790), "固件按当前页定向刷新控件，页面切换后会用缓存数据立即补齐显示。", font=load_font(27), fill=GRAY)
    return save(img, "04_serial_screen_pages.png")


def draw_mini_program_flow() -> Path:
    img, draw = new_canvas()
    title(draw, "小程序交互流程")
    # Phone mockup
    rounded_box(draw, (100, 150, 500, 760), "#FFFFFF", GREEN, 5, 38)
    draw.text((185, 185), "空气管家", font=load_font(34, True), fill=GREEN)
    for i, (head, val) in enumerate([("温度", "26.3 C"), ("湿度", "58 %"), ("空气", "75")]):
        x = 135 + i * 115
        rounded_box(draw, (x, 260, x + 95, 345), "#F8FAF7", LINE, 2, 12)
        draw.text((x + 18, 278), val, font=load_font(22, True), fill=ORANGE)
        draw.text((x + 25, 312), head, font=load_font(18), fill=GRAY)
    for i, name in enumerate(["净化", "新风", "加湿", "自动", "节能"]):
        y = 395 + i * 58
        draw.text((145, y), name, font=load_font(23), fill=DARK)
        for j in range(4):
            fill = ORANGE if j == (i % 4) else "#EEF2EE"
            draw.rounded_rectangle((260 + j * 46, y - 4, 294 + j * 46, y + 30), radius=8, fill=fill)
            draw.text((270 + j * 46, y + 1), str(j), font=load_font(17, True), fill="#FFFFFF" if fill == ORANGE else GREEN)
    steps = [
        ("1 填写 ESP32 地址", 665, 200),
        ("2 读取 /api/state", 665, 330),
        ("3 控制 /api/device", 665, 460),
        ("4 查看 /api/history", 665, 590),
    ]
    for txt, x, y in steps:
        rounded_box(draw, (x, y, x + 560, y + 82), "#FFFFFF", LINE, 3, 14)
        draw.text((x + 28, y + 23), txt, font=load_font(27, True), fill=DARK)
        if y < 590:
            arrow(draw, (945, y + 82), (945, y + 126), GREEN, 4)
    arrow(draw, (500, 455), (665, 455), ORANGE, 5)
    draw.text((100, 810), "小程序不另建云服务器，适合答辩现场和同一 WiFi 下的手机端控制。", font=load_font(27), fill=GRAY)
    return save(img, "05_mini_program_flow.png")


def draw_mcp_bridge_flow() -> Path:
    img, draw = new_canvas()
    title(draw, "AI/MCP 语音控制链路")
    xs = [70, 370, 690, 1010, 1310]
    labels = [
        ("语音指令", "打开净化器二档"),
        ("小智平台", "MCP 接入点"),
        ("本地桥接", "smart_home_bridge.py"),
        ("ESP32 API", "/api/device"),
        ("执行器动作", "GPIO13/14/21"),
    ]
    for i, ((head, desc), x) in enumerate(zip(labels, xs)):
        rounded_box(draw, (x, 280, x + 220, 420), "#FFFFFF", GREEN if i in [0, 4] else LINE, 4 if i in [0, 4] else 2, 20)
        draw_wrapped(draw, (x + 18, 310), head, load_font(27, True), GREEN, 184, center=True)
        draw_wrapped(draw, (x + 18, 358), desc, load_font(20), GRAY, 184, center=True)
        if i < len(xs) - 1:
            arrow(draw, (x + 220, 350), (xs[i + 1], 350), ORANGE, 5)
    rounded_box(draw, (360, 535, 1240, 665), "#FFFFFF", LINE, 2, 18)
    draw_wrapped(draw, (400, 565), "语音控制没有绕过本地控制逻辑，而是复用同一组 HTTP API 和设备状态。这样屏幕、小程序、AI 控制看到的是同一个结果。", load_font(27), DARK, 800)
    return save(img, "06_mcp_bridge_flow.png")


def draw_auto_eco_logic() -> Path:
    img, draw = new_canvas()
    title(draw, "自动模式与节能模式逻辑")
    rounded_box(draw, (90, 165, 450, 260), "#FFFFFF", GREEN, 4, 18)
    draw.text((132, 195), "读取环境状态", font=load_font(30, True), fill=GREEN)
    decision = [(720, 130, "湿度 < 40%", "加湿二档"), (720, 315, "空气评分偏低", "净化 + 新风"), (720, 500, "温度偏高", "提高新风档位")]
    arrow(draw, (450, 212), (640, 212), ORANGE, 5)
    for x, y, cond, action in decision:
        draw.polygon([(x, y), (x + 210, y + 70), (x, y + 140), (x - 210, y + 70)], fill="#FFFFFF", outline=GREEN)
        draw_wrapped(draw, (x - 120, y + 45), cond, load_font(24, True), DARK, 240, center=True)
        rounded_box(draw, (1030, y + 30, 1390, y + 110), "#FFFFFF", LINE, 3, 16)
        draw.text((1070, y + 52), action, font=load_font(27, True), fill=ORANGE)
        arrow(draw, (x + 210, y + 70), (1030, y + 70), ORANGE, 4)
    rounded_box(draw, (90, 555, 520, 705), "#FFFFFF", ORANGE, 4, 18)
    draw_wrapped(draw, (130, 590), "节能模式：关闭自动模式，并把净化、新风、加湿档位全部置 0。", load_font(26, True), DARK, 350)
    draw.text((95, 790), "自动模式负责按规则调节，节能模式负责一键收束功耗，两者互斥。", font=load_font(27), fill=GRAY)
    return save(img, "07_auto_eco_logic.png")


def draw_test_workflow() -> Path:
    img, draw = new_canvas()
    title(draw, "测试验证流程")
    items = ["传感器读数", "串口屏显示", "触摸控制", "小程序 API", "AI/MCP", "自动联动", "连续运行"]
    x = 110
    y = 230
    for i, item in enumerate(items):
        w = 185
        h = 120
        fill = "#FFFFFF" if i % 2 == 0 else "#F9FBF9"
        rounded_box(draw, (x, y, x + w, y + h), fill, GREEN if i == 0 else LINE, 3, 18)
        draw.text((x + 24, y + 20), f"{i+1}", font=load_font(36, True), fill=ORANGE)
        draw_wrapped(draw, (x + 58, y + 38), item, load_font(23, True), DARK, 105, center=True)
        if i < len(items) - 1:
            arrow(draw, (x + w, y + h // 2), (x + w + 55, y + h // 2), GREEN, 4)
        x += w + 70
        if i == 3:
            x = 260
            y = 515
    rounded_box(draw, (1120, 560, 1450, 705), "#FFFFFF", ORANGE, 4, 18)
    draw_wrapped(draw, (1150, 590), "验证重点：三种入口控制同一套状态，刷新不丢数据，执行器动作和界面状态一致。", load_font(25), DARK, 270)
    return save(img, "08_test_workflow.png")


def draw_application_scenarios() -> Path:
    img, draw = new_canvas()
    title(draw, "应用场景")
    cx, cy = 800, 440
    draw.ellipse((610, 250, 990, 630), fill="#FFFFFF", outline=GREEN, width=6)
    draw_wrapped(draw, (682, 360), "室内环境\n智能调节", load_font(36, True), GREEN, 240, center=True)
    scenes = [
        ("家庭卧室", "夜间温湿度和空气提醒", 180, 180),
        ("宿舍", "多人空间通风管理", 1135, 180),
        ("教室", "人员密集时空气看板", 165, 585),
        ("小办公室", "手机和语音快速控制", 1120, 585),
    ]
    for head, desc, x, y in scenes:
        rounded_box(draw, (x, y, x + 330, y + 150), "#FFFFFF", LINE, 3, 18)
        draw.text((x + 28, y + 28), head, font=load_font(29, True), fill=GREEN)
        draw_wrapped(draw, (x + 28, y + 76), desc, load_font(23), DARK, 270)
        arrow(draw, (cx, cy), (x + 165, y + 75), ORANGE, 4)
    draw.text((115, 790), "价值落点不是堆功能，而是把“看见环境、及时控制、自动调节”放到真实小空间里。", font=load_font(27), fill=GRAY)
    return save(img, "09_application_scenarios.png")


def draw_screen_miniapp_combo() -> Path:
    img, draw = new_canvas()
    title(draw, "三种交互入口")
    cols = [
        ("串口屏", "现场查看和触摸控制", "不用手机，站在设备旁就能看状态。"),
        ("小程序", "手机端读取和控制", "同一 WiFi 下查看历史、切换档位。"),
        ("AI 语音", "口语指令变成工具调用", "把“打开净化器二档”转成 API 控制。"),
    ]
    for i, (head, sub, desc) in enumerate(cols):
        x = 110 + i * 500
        rounded_box(draw, (x, 210, x + 390, 610), "#FFFFFF", GREEN if i == 1 else LINE, 4 if i == 1 else 2, 24)
        draw.text((x + 42, 250), head, font=load_font(36, True), fill=GREEN)
        draw.text((x + 42, 315), sub, font=load_font(25, True), fill=ORANGE)
        draw_wrapped(draw, (x + 42, 390), desc, load_font(24), DARK, 310)
        draw.ellipse((x + 140, 500, x + 250, 610), fill="#F8FAF7", outline=LINE, width=2)
    draw.text((115, 760), "三个入口的结果都会回写到同一份设备状态，答辩演示时可以任选其中一种入口开始。", font=load_font(27), fill=GRAY)
    return save(img, "10_interaction_channels.png")


def create_figures() -> dict[str, Figure]:
    specs = [
        ("system_architecture", "图 1 系统总体架构图", draw_system_architecture),
        ("data_control_flow", "图 2 数据流与控制流图", draw_data_control_flow),
        ("hardware_wiring", "图 3 硬件连接示意图", draw_hardware_wiring),
        ("serial_screen_pages", "图 4 串口屏页面结构图", draw_serial_screen_pages),
        ("mini_program_flow", "图 5 小程序交互流程图", draw_mini_program_flow),
        ("mcp_bridge_flow", "图 6 AI/MCP 语音控制链路图", draw_mcp_bridge_flow),
        ("auto_eco_logic", "图 7 自动模式与节能模式逻辑图", draw_auto_eco_logic),
        ("test_workflow", "图 8 测试验证流程图", draw_test_workflow),
        ("application_scenarios", "图 9 应用场景图", draw_application_scenarios),
        ("interaction_channels", "图 10 三种交互入口说明图", draw_screen_miniapp_combo),
    ]
    figures: dict[str, Figure] = {}
    for key, caption, fn in specs:
        figures[key] = Figure(key, caption, fn())
    return figures


def gpio_table() -> TableData:
    return TableData(
        "表 1 主要 GPIO 与接口分配",
        ["模块", "接口/引脚", "作用"],
        [
            ["DHT11", "GPIO18", "读取温度和湿度"],
            ["MQ135", "GPIO1 / ADC1_CH0", "读取空气质量模拟量"],
            ["TJC 串口屏", "UART2 GPIO41 TX / GPIO42 RX", "本地显示和触摸事件"],
            ["净化指示", "GPIO13 / LEDC PWM", "红色 LED 亮度表示净化档位"],
            ["加湿指示", "GPIO14 / LEDC PWM", "蓝色 LED 亮度表示加湿档位"],
            ["新风扇叶", "GPIO21 / 50Hz PWM", "舵机在 0 到 180 度范围按档位往复摆动"],
        ],
    )


def api_table() -> TableData:
    return TableData(
        "表 2 小程序 HTTP API",
        ["接口", "方法", "用途"],
        [
            ["/api/state", "GET", "读取温湿度、空气评分、设备档位、模式状态"],
            ["/api/history", "GET", "读取最近 30 条空气质量历史数据"],
            ["/api/device", "POST", "控制净化、新风、加湿的开关和档位"],
            ["/api/mode", "POST", "控制自动模式和节能模式"],
            ["/api/environment", "POST", "手动输入环境数据或切换舒适、高温、干燥、污染场景"],
        ],
    )


def mcp_table() -> TableData:
    return TableData(
        "表 3 AI/MCP 工具列表",
        ["工具名", "作用", "示例口语"],
        [
            ["home_get_state", "读取当前环境和设备状态", "现在空气怎么样"],
            ["home_set_purifier", "控制净化 LED 档位", "打开净化器二档"],
            ["home_set_fresh_air", "控制新风舵机扇叶档位", "开启新风三档"],
            ["home_set_humidifier", "控制加湿 LED 档位", "关闭加湿器"],
            ["home_set_auto", "开启或关闭自动模式", "打开自动模式"],
            ["home_set_eco", "开启或关闭节能模式", "进入节能模式"],
            ["home_set_environment_preset", "切换手动环境场景", "模拟污染环境"],
            ["home_get_advice", "读取舒适度和建议", "现在需要开新风吗"],
        ],
    )


def test_table() -> TableData:
    return TableData(
        "表 4 功能验证记录",
        ["测试项", "操作方法", "通过现象"],
        [
            ["传感器读取", "上电后观察串口日志和屏幕首页", "温湿度、空气评分能随采样周期刷新"],
            ["串口屏显示", "停留首页和空气详情页观察刷新", "上一次数据保持在屏幕上，新数据到达后覆盖更新"],
            ["串口屏控制", "点击净化、新风、加湿、自动、节能热区", "对应档位和模式状态变化"],
            ["小程序控制", "访问 ESP32 地址并调用设备控制按钮", "HTTP 返回状态变化，外设动作同步"],
            ["AI 语音控制", "通过 MCP 桥接发出控制工具调用", "ESP32 API 接收到请求并执行动作"],
            ["自动模式", "切换干燥或污染场景", "加湿、净化、新风档位按规则调整"],
            ["节能模式", "开启节能", "自动模式关闭，三个执行器档位归零"],
        ],
    )


def build_chapters() -> tuple[list[str], list[Chapter], list[TableData]]:
    abstract = [
        "本作品面向家庭、宿舍、教室和小型办公室等室内空间，设计并实现了一套基于 ESP32S3 的室内环境智能调节系统。系统围绕温湿度和空气质量两个最直观的环境指标展开，通过 DHT11 采集温湿度，通过 MQ135 采集空气质量模拟量，再由 ESP32S3 完成状态计算、显示刷新和执行器控制。",
        "作品的交互方式分为三类。第一类是 TJC 串口屏，负责现场显示温湿度、空气评分、环境建议和设备状态，也支持触摸控制净化、新风、加湿、自动和节能模式。第二类是微信小程序，用户在同一局域网中填写 ESP32 地址后，可以读取当前状态、查看近期历史、切换设备档位。第三类是 AI 语音控制，系统通过 MCP 工具和本地桥接程序把口语指令转成 ESP32 HTTP API 请求，让语音入口也能控制同一套设备状态。",
        "执行器部分采用低成本、容易观察的方式完成联动展示：红色 LED 表示净化，蓝色 LED 表示加湿，二者通过 PWM 亮度表示不同档位；新风部分使用 GPIO21 输出舵机控制信号，带动扇叶在 0 到 180 度范围内按档位往复摆动。自动模式根据温度、湿度和空气评分调整设备档位，节能模式则关闭自动模式和所有执行器。",
        "从应用角度看，本作品不是单独的传感器读数页面，而是把环境感知、现场显示、手机控制、语音控制和自动联动放在一个小型系统里。它适合用于室内空气状态提醒、宿舍通风管理、教室空气看板和小办公室环境管理，后续可以继续扩展云端记录、更多空气传感器和真实家电控制模块。",
    ]

    chapters = [
        Chapter(
            "第 1 章 作品概述",
            [
                Section(
                    "1.1 作品背景与国内外研究现状",
                    [
                        "室内环境质量和人的日常生活联系很紧。家庭、宿舍、教室和办公室里，温度、湿度、通风状态和空气质量都会影响体感舒适度。现实使用中，很多人只能靠感觉判断环境变化，等到觉得闷、干、热或者空气不好时才去开窗、开风扇或开加湿设备。这样的处理方式不够及时，也不方便做连续观察。",
                        "市面上常见的温湿度计可以显示数据，但一般只负责显示，不能直接控制设备；空气质量检测仪能提供空气状态提醒，但和净化器、风扇、加湿器之间往往还需要用户手动切换；普通智能家居设备可以远程开关，但很多设备缺少本地环境判断，只是把手机变成遥控器。对于小型空间来说，真正需要的是一套成本适中、显示直观、控制入口统一的系统。",
                        "物联网方向的发展思路已经比较清楚：传感器负责感知，主控负责判断，屏幕和手机负责交互，执行器负责动作。我们在作品中把这个思路缩小到一个可以上电运行的室内环境系统里。系统不追求堆很多昂贵硬件，而是把现有的 ESP32S3、串口屏、DHT11、MQ135、LED 和舵机组合起来，先把从数据到控制的闭环跑通。",
                    ],
                    figure_key="application_scenarios",
                ),
                Section(
                    "1.2 应用领域",
                    [
                        "本作品主要面向小型室内空间。家庭卧室中，用户可以通过串口屏或手机查看温湿度和空气评分，了解晚上是否需要加湿或通风；宿舍中，多个成员共用空间，空气状态不容易统一感知，系统可以作为一个放在桌面上的环境提醒终端；教室和小办公室中，人员较多时空气变化快，本作品可以作为低成本空气看板和控制入口。",
                        "在这些场景里，用户不一定需要复杂的云平台，但需要能看见实时数据、能快速操作、能让系统在明显异常时自动做出动作。作品的三种入口正好对应不同使用习惯：人在设备旁边时看串口屏，离设备有一段距离时用小程序，不方便操作时用语音。",
                    ],
                ),
                Section(
                    "1.3 主要功能与特性",
                    [
                        "系统已经完成温湿度检测、空气质量评分、串口屏显示、小程序读取与控制、AI/MCP 语音控制、净化/新风/加湿执行器控制、自动模式、节能模式和近期空气质量历史记录。串口屏首页负责快速查看状态，空气详情页负责查看评分和曲线，智能家居页负责直接控制三个执行器和两个模式。",
                        "净化功能通过 GPIO13 的红色 LED 表示，档位越高亮度越高；加湿功能通过 GPIO14 的蓝色 LED 表示，档位同样由亮度区分；新风功能通过 GPIO21 舵机带动扇叶摆动，档位越高摆动范围和节奏越明显。这样的设计虽然结构简单，但非常适合现场答辩观察，因为每个控制结果都能直接看见。",
                        "小程序通过局域网 HTTP API 访问 ESP32S3，不需要额外云服务器。用户填写 ESP32 的 IP 和端口后，可以读取状态、切换设备档位、开启自动或节能模式，也可以输入手动环境数据来验证自动控制逻辑。AI 语音控制则通过 MCP 桥接服务复用同一组 API，避免出现语音控制和手机控制状态不一致的问题。",
                    ],
                    figure_key="interaction_channels",
                ),
                Section(
                    "1.4 创新性说明",
                    [
                        "本作品的创新点不在于单个传感器或某一个按钮，而在于把几个常见入口接到同一套控制状态上。串口屏、小程序和 AI 语音都不是各做各的，而是共同读写 ESP32S3 内部的设备状态。这样用户从屏幕操作后，小程序能读到变化；从语音入口控制后，串口屏也能继续显示同一套结果。",
                        "另一个特点是自动模式和手动环境场景结合。真实环境不一定刚好出现高温、干燥或污染状态，为了让联动逻辑能在答辩现场稳定展示，系统提供舒适、高温、干燥、污染几种场景。这样可以清楚说明系统如何从环境状态推导到净化、新风、加湿的动作。",
                    ],
                ),
                Section(
                    "1.5 设计流程",
                    [
                        "设计流程从硬件连接开始，先确认 ESP32S3 可读取 DHT11 和 MQ135，再把数据稳定显示到串口屏。随后补充智能家居控制器，把净化、新风、加湿三个动作统一成 0 到 3 档状态。完成本地控制后，再加入 HTTP API 和小程序，让手机端能够读取和修改同一套状态。最后接入 MCP 桥接，把语音入口也接到 HTTP API 上。",
                        "调试过程中，串口屏显示连续性和温湿度数据稳定性是重点。系统最终采用缓存最新数据、按当前页面定向刷新、页面切换后立即补齐控件的方式处理屏幕显示；传感器读数异常时显示占位符，避免把错误湿度或温度直接展示给用户。",
                    ],
                    figure_key="data_control_flow",
                ),
            ],
        ),
        Chapter(
            "第 2 章 需求分析",
            [
                Section(
                    "2.1 开发背景与目的",
                    [
                        "作品开发的直接目的，是做出一个能在现场真实运行的室内环境调节系统。它需要读到环境数据，也需要把数据用普通用户能看懂的方式展示出来，还要能控制一些可观察的执行器。只显示数字不够，只做手机按钮也不够，系统要能把感知和动作接起来。",
                        "从用户角度看，室内空气管理经常遇到三个问题：第一，环境变化不直观，用户不知道现在是偏干、偏闷还是空气较差；第二，控制入口分散，屏幕、手机和语音设备之间常常没有统一状态；第三，设备开关依赖人工，用户容易忘记关闭或不知道何时调整档位。本作品围绕这三个问题展开设计。",
                    ],
                ),
                Section(
                    "2.2 用户群体分析",
                    [
                        "家庭用户更关心操作是否简单。老人或孩子不一定愿意打开复杂 App，因此串口屏首页必须能直接显示温度、湿度和空气状态；年轻用户更习惯手机控制，因此小程序要能快速连接和控制设备；语音入口适合做一些不方便手动操作的场景，例如睡前说一句打开加湿或开启节能。",
                        "宿舍、教室和办公室用户更关心公共空间的状态展示。设备放在桌面或讲台上时，串口屏可以作为一个小看板；管理员或同学也可以用手机查看状态。对于这类空间，系统的价值不只是打开某个设备，而是让空气状态变成人人看得见、能讨论、能及时处理的信息。",
                    ],
                ),
                Section(
                    "2.3 同类方案分析",
                    [
                        "传统温湿度计价格低、使用简单，但功能停留在显示层，不能直接控制设备。空气质量检测仪比温湿度计多了空气状态判断，但很多产品和执行器是分开的，用户看完还要手动打开净化器或风扇。普通智能开关可以远程控制电器，但如果没有环境数据参与，它只是一个远程按钮。",
                        "本作品采用的方案介于这些产品之间。它不是昂贵的完整家居系统，也不是单个传感器模块，而是一套小型闭环：传感器给出数据，ESP32S3 做状态整理，屏幕和小程序显示状态，执行器根据手动或自动控制做出动作，AI 语音作为额外入口接入。",
                    ],
                    table=TableData(
                        "表 5 同类方案对比",
                        ["方案", "优势", "不足", "本作品对应改进"],
                        [
                            ["温湿度计", "便宜、直观", "只能显示，不能联动", "加入净化、新风、加湿控制"],
                            ["空气检测仪", "能提示空气状态", "控制动作仍需人工", "空气评分可触发自动调节"],
                            ["智能开关", "手机端控制方便", "缺少本地环境判断", "传感数据和设备状态统一"],
                            ["完整智能家居", "功能丰富", "成本和部署复杂度较高", "采用 ESP32S3 做低成本版本"],
                        ],
                    ),
                ),
                Section(
                    "2.4 主要功能需求",
                    [
                        "系统功能需求可以分为五类。第一类是环境采集，要能获取温湿度和空气质量数据；第二类是本地显示，要能在串口屏上持续显示上一次有效数据，并在新数据到达时覆盖更新；第三类是设备控制，要能控制净化、新风、加湿三个执行器，并用档位区分强弱；第四类是远程交互，要能通过小程序读取状态和控制设备；第五类是语音交互，要能通过 MCP 工具把口语指令转为控制动作。",
                        "自动和节能是两个关键模式。自动模式要根据环境状态调整设备，而不是只改变一个页面图标；节能模式要明确关闭自动模式和所有执行器，给用户一个一键收束的入口。两个模式互斥，避免系统一边自动打开设备，一边又处在节能状态。",
                    ],
                ),
                Section(
                    "2.5 主要性能指标响应",
                    [
                        "本作品的性能指标更偏向可用性和一致性。显示方面，串口屏应保持上一次数据，不应只在刷新瞬间显示；控制方面，串口屏、小程序和 AI 语音操作后，设备状态要保持一致；历史方面，系统保存最近 30 条采样数据，用于展示近期空气评分变化；联动方面，手动环境场景切换后，自动模式能在采样更新中调整设备档位。",
                        "考虑到当前硬件以 DHT11 和 MQ135 为主，报告中的数据指标不写成高精度工业测量。当前版本重点证明系统结构和控制链路可用，后续可以把传感器扩展为 SHT30、PM2.5、CO2、TVOC 等更细分模块，再提高检测精度和指标覆盖范围。",
                    ],
                ),
            ],
        ),
        Chapter(
            "第 3 章 技术方案",
            [
                Section(
                    "3.1 系统整体架构介绍",
                    [
                        "系统整体由感知与执行层、ESP32S3 主控层、交互层和应用层组成。感知与执行层包含 DHT11、MQ135、红色 LED、蓝色 LED 和舵机扇叶；主控层负责采样、评分、状态缓存、自动规则和 HTTP 服务；交互层包含 TJC 串口屏、微信小程序和 AI/MCP 桥接；应用层对应家庭、宿舍、教室和办公室等真实使用场景。",
                        "这样的分层有一个好处：每一层的职责清楚。传感器不负责显示，屏幕不负责判断规则，AI 语音不直接操作 GPIO，而是通过统一接口交给 ESP32S3 控制。后续如果更换传感器或接入真实净化器，只需要替换对应层的模块，不需要推倒整套交互逻辑。",
                    ],
                    figure_key="system_architecture",
                ),
                Section(
                    "3.2 系统流程介绍",
                    [
                        "系统上电后，ESP32S3 初始化串口屏、传感器、PWM 通道和 HTTP 服务。传感器任务按周期读取 DHT11 和 MQ135 数据，主控把温度、湿度、空气原始值和空气评分整理成一份环境状态。串口屏刷新时从这份状态取数据，小程序访问 `/api/state` 时也读取同一份状态。",
                        "当用户从任意入口发起控制时，动作都会进入智能家居控制器。串口屏触摸事件会解析成 `BTN,DEVICE` 或 `BTN,MODE` 指令；小程序会发 HTTP POST 请求；AI 语音通过 MCP 桥接程序转发到 HTTP API。控制器更新状态后再输出 PWM 或模式变化，因此三个入口不会各自维护一份状态。",
                    ],
                    figure_key="data_control_flow",
                ),
                Section(
                    "3.3 硬件系统介绍",
                    [
                        "硬件以 ESP32S3 为核心。DHT11 接 GPIO18，用于读取温度和湿度；MQ135 接 GPIO1 的 ADC1_CH0，用于读取空气质量模拟量；TJC 串口屏使用 UART2，ESP32S3 的 GPIO41 作为 TX 连接屏幕 RX，GPIO42 作为 RX 连接屏幕 TX，波特率为 9600。",
                        "执行器部分采用便于观察的形式。GPIO13 输出 PWM 控制红色 LED，表示净化器档位；GPIO14 输出 PWM 控制蓝色 LED，表示加湿器档位；GPIO21 输出 50Hz PWM 控制舵机，舵机带动扇叶做新风动作。档位用 0、1、2、3 表示，0 为关闭，1 到 3 表示强度逐步提高。",
                    ],
                    figure_key="hardware_wiring",
                    table=gpio_table(),
                ),
                Section(
                    "3.4 电路与接口模块说明",
                    [
                        "DHT11 使用单总线方式通信，对时序要求比较严格。固件中对异常数据做了校验，湿度超出范围或校验和错误时不会把错误值直接显示出来，而是显示占位符。这样可以避免上电调试时出现湿度大于 100% 后仍然被当成有效结果展示。",
                        "MQ135 通过 ADC 读取原始值，再转换为空气评分和空气等级。当前版本不把 MQ135 原始值直接写成某一种气体浓度，而是把它作为空气状态变化的输入，这样表述更贴近当前硬件实际。后续如果加入标定流程和更多传感器，可以再把评分拆成 CO2、TVOC 或 PM2.5 等具体指标。",
                        "串口屏使用 TJC 指令协议，发送命令时以 `FF FF FF` 结尾。固件按当前页面刷新控件，避免给不存在于当前页的控件持续发命令。页面切换后，固件会用缓存的最新数据立即刷新当前页，从而让显示更连续。",
                    ],
                ),
                Section(
                    "3.5 软件系统介绍",
                    [
                        "软件上，智能家居控制器负责维护 `purifier_level`、`fresh_air_level`、`humidifier_level`、`auto_mode`、`eco_mode` 等状态。每次状态变化后，控制器会调用对应输出函数，把状态落到 LED PWM 或舵机 PWM 上。这样报告和答辩中讲到的“档位”，在固件里有明确的数据字段和输出动作。",
                        "HTTP 服务运行在 ESP32S3 局域网地址的 8080 端口，提供状态查询、历史查询、设备控制、模式控制和环境场景输入。HTTP API 的作用不是单独服务小程序，它也给 AI/MCP 桥接层复用。这样手机端和语音端的控制路径最终汇入同一个控制器。",
                    ],
                    table=api_table(),
                ),
                Section(
                    "3.6 AI/MCP 与小程序交互方案",
                    [
                        "小程序方案采用局域网直连。用户在页面中输入 ESP32S3 的地址，例如 `192.168.1.23:8080`，小程序先调用 `/api/state` 和 `/api/history` 刷新页面，再通过按钮调用 `/api/device` 和 `/api/mode`。它适合现场答辩和同一 WiFi 下的实际控制。",
                        "AI 语音控制通过外部 Python 桥接服务完成。小智平台的 MCP 接入点收到工具调用后，本地 `smart_home_bridge.py` 把它转换为 HTTP 请求发给 ESP32S3。比如“打开净化器二档”会调用 `home_set_purifier(power=true, level=2)`，桥接层再请求 `/api/device`，最后由 ESP32S3 改变 GPIO13 的 PWM 亮度。",
                    ],
                    figure_key="mcp_bridge_flow",
                    table=mcp_table(),
                ),
            ],
        ),
        Chapter(
            "第 4 章 方案实现",
            [
                Section(
                    "4.1 整体实现说明",
                    [
                        "当前工程已经把传感器、串口屏、执行器控制、小程序 API 和 AI/MCP 桥接串成了一条完整链路。系统运行时，传感器数据每隔一段时间刷新一次，最新状态会进入缓存；屏幕、小程序和语音入口都围绕这份状态工作。这样做的目的，是让作品在答辩现场能从多个角度展示同一套系统，而不是几个互不相干的小模块。",
                        "代码结构上，`serial_hmi.cc/.h` 负责串口屏显示和事件解析，`smart_home_controller.cc/.h` 负责执行器状态、自动和节能模式、MCP 工具注册，`smart_home_http_server.cc/.h` 负责 HTTP API，小程序工程放在 `docs/mini_program_demo`，MCP 桥接服务放在 `tools/xiaozhi_mcp_bridge`。",
                    ],
                ),
                Section(
                    "4.2 串口屏实现",
                    [
                        "串口屏分为首页、空气详情、智能家居控制、AI 与设置页面。首页显示温度、湿度、空气状态、空气评分和建议，适合快速查看；空气详情页显示 MQ135 原始值、空气评分和近期曲线，适合解释传感器数据；智能家居页负责净化、新风、加湿、自动和节能操作；AI 与设置页负责显示 AI 状态和手动环境场景。",
                        "屏幕刷新采用缓存最新数据的方式。传感器读数更新时，固件先记录最新状态，再根据当前页面刷新对应控件。如果停留在首页，首页数据会保持在屏幕上；如果切到空气详情页，固件会把最近的空气评分曲线重新写入曲线控件。这样可以改善只在刷新瞬间显示数据的问题。",
                    ],
                    figure_key="serial_screen_pages",
                ),
                Section(
                    "4.3 小程序实现",
                    [
                        "小程序页面包含 ESP32 地址输入、环境状态、手动输入数据、模式控制、设备档位控制和近期空气质量记录。用户点击连接后，小程序会保存地址并读取状态。设备控制区用 0、1、2、3 四个按钮表示关闭和三档，点击后立刻调用 HTTP API。",
                        "手动输入数据区用于现场验证自动模式。用户可以输入温度、湿度和空气评分，也可以点击舒适、高温、干燥、污染几个预设场景。这个设计让自动逻辑可以在没有真实污染或干燥环境的情况下被清楚展示，评委能看到系统从环境状态到设备动作的变化。",
                    ],
                    figure_key="mini_program_flow",
                ),
                Section(
                    "4.4 智能家居执行器实现",
                    [
                        "净化、加湿、新风三个执行器采用统一的档位模型。净化档位写入 `purifier_level`，输出到 GPIO13 红色 LED；加湿档位写入 `humidifier_level`，输出到 GPIO14 蓝色 LED；新风档位写入 `fresh_air_level`，输出到 GPIO21 舵机。三个设备都支持 0 到 3 档，便于屏幕、小程序和 AI 统一表达。",
                        "LED 通过 LEDC PWM 输出亮度。档位 0 为熄灭，档位 1 到 3 逐步增加亮度。舵机采用 50Hz PWM，固件按档位设置摆动范围和步进速度，在 0 到 180 度范围内往复摆动。实物上，扇叶插在舵机上，摆动效果可以直接代表新风档位。",
                    ],
                ),
                Section(
                    "4.5 自动模式和节能模式实现",
                    [
                        "自动模式开启后，系统会根据最近一次环境状态调整设备。湿度低于阈值时提高加湿档位；空气评分偏低时提高净化和新风档位；温度偏高时提高新风档位。自动模式不是单独的页面状态，它会真实改写执行器档位。",
                        "节能模式开启后，系统会关闭自动模式，并把净化、新风、加湿三个档位全部置 0。这个设计让用户在离开房间、夜间安静时段或需要降低功耗时，可以一键把设备收束到关闭状态。自动和节能互斥，避免出现逻辑冲突。",
                    ],
                    figure_key="auto_eco_logic",
                ),
                Section(
                    "4.6 AI 语音控制实现",
                    [
                        "固件内部注册了 `self.home.*` 工具，用于描述智能家居能力；外部桥接服务则通过 FastMCP 暴露 `home_get_state`、`home_set_purifier`、`home_set_fresh_air`、`home_set_humidifier`、`home_set_auto`、`home_set_eco` 等工具。语音平台只需要选择合适工具，真正的设备控制仍由 ESP32S3 完成。",
                        "这种方式的优点是边界清楚。AI 不需要知道 GPIO13 是红色 LED，也不需要知道舵机脉宽怎么写，它只负责把“开新风三档”理解为工具调用。桥接层负责把工具参数变成 HTTP JSON，ESP32S3 负责执行动作并返回最新状态。",
                    ],
                    figure_key="mcp_bridge_flow",
                ),
                Section(
                    "4.7 工程成果展示",
                    [
                        "目前实物部分已经具备可观察的硬件输出，串口屏能够显示首页、空气详情和 AI 设置页面，小程序工程可以通过局域网访问 ESP32S3，AI/MCP 桥接服务可以把工具调用转发到 HTTP API。报告中的实物图和屏幕截图放在本节，便于读者把前面的技术方案和实际作品对应起来。",
                        "工程成果的重点不是外观包装，而是功能链路完整。传感器有数据来源，屏幕有持续显示，小程序有接口，AI 有工具链路，执行器有可观察动作，自动和节能有明确规则。答辩时可以按这条顺序展示，逻辑会比较顺。",
                    ],
                ),
            ],
        ),
        Chapter(
            "第 5 章 测试报告",
            [
                Section(
                    "5.1 测试环境与定量标准",
                    [
                        "测试环境包括 ESP32S3 主控板、DHT11 温湿度传感器、MQ135 空气质量传感器、TJC 串口屏、红色 LED、蓝色 LED、舵机扇叶、运行小程序的电脑或手机，以及运行 MCP 桥接服务的本地电脑。测试时 ESP32S3、小程序设备和桥接电脑需要处在同一局域网中。",
                        "本作品的测试标准以功能通过为主。传感器读数要能进入系统状态；屏幕要能持续显示；控制入口要能改变设备状态；执行器动作要和状态一致；自动模式要能根据环境场景调整；节能模式要能关闭执行器。对于 DHT11 和 MQ135，当前版本不写高精度测量结论，而是验证数据读取和控制链路。",
                    ],
                    figure_key="test_workflow",
                ),
                Section(
                    "5.2 传感器数据显示测试",
                    [
                        "上电后，通过串口 monitor 观察 DHT11 和 MQ135 的读取结果。DHT11 正常时输出温度和湿度，异常时固件记录错误并在显示侧使用占位符，避免错误值进入界面；MQ135 正常时输出 ADC 原始值和空气等级。经过前期排查，湿度大于 100% 的异常读数不会再作为正常数据展示。",
                        "屏幕首页应显示温度、湿度、空气状态和建议；空气详情页应显示空气评分、MQ135 原始值和曲线。小程序 `/api/state` 返回 JSON 中应包含温湿度、空气评分、设备档位和模式状态。这样可以从串口、屏幕、小程序三个位置确认数据进入系统。",
                    ],
                ),
                Section(
                    "5.3 串口屏显示和触摸控制测试",
                    [
                        "显示测试的重点是连续性。停留在一个页面时，数据应保持在屏幕上，直到下一次采样到达后覆盖更新，不应在两次刷新之间回到初始状态。切换页面后，固件应使用缓存数据刷新当前页控件，避免页面只有背景没有数据。",
                        "触摸控制测试包括点击净化、新风、加湿、自动和节能热区。每次点击后，串口日志应解析出对应事件，控制器状态应变化，外设应出现对应动作。页面快速切换时，固件有防抖处理，避免短时间连续跳页造成画面不稳定。",
                    ],
                ),
                Section(
                    "5.4 小程序 HTTP API 测试",
                    [
                        "小程序测试先从浏览器访问 `http://<ESP32_IP>:8080/api/state` 开始，确认 HTTP 服务已经启动并返回 JSON。随后在小程序中填写 `<ESP32_IP>:8080`，点击连接，观察页面是否显示环境状态、模式按钮、设备档位按钮和近期空气质量记录。",
                        "设备控制测试分别点击净化、新风、加湿的 1、2、3、0 档，观察 `/api/state` 中的 `purifier_level`、`fresh_air_level`、`humidifier_level` 是否变化，并观察 GPIO13 红色 LED、GPIO14 蓝色 LED 和 GPIO21 舵机扇叶动作是否同步。",
                    ],
                    table=test_table(),
                ),
                Section(
                    "5.5 AI/MCP 语音控制测试",
                    [
                        "AI/MCP 测试先确认 ESP32 HTTP API 可访问，再启动本地 MCP 桥接服务。桥接服务需要设置 `ESP32_BASE_URL`，指向 `http://<ESP32_IP>:8080`，再通过小智平台 MCP 接入点连接。连接成功后，语音平台可以看到 `home_*` 工具列表。",
                        "测试口令可以选择“打开净化器二档”“开启新风三档”“关闭加湿器”“打开自动模式”“进入节能模式”“模拟污染环境”。每条口令对应一个工具调用，工具调用返回的 JSON 状态应和小程序读取结果一致，外设动作也应同步变化。",
                    ],
                ),
                Section(
                    "5.6 自动/节能模式联动测试",
                    [
                        "自动模式测试建议先打开自动模式，再通过小程序或串口屏切换手动环境场景。选择干燥场景时，加湿档位应提高；选择污染场景时，净化和新风档位应提高；选择舒适场景时，设备档位应回到较低状态。这样可以直观看到规则起作用。",
                        "节能模式测试则先让任意执行器处于非 0 档，再开启节能。通过 `/api/state` 和实物观察确认 `eco_mode` 为 true，`auto_mode` 为 false，三个执行器档位都变为 0。再次关闭节能后，系统保留手动控制入口，但不会自动恢复之前的高档位。",
                    ],
                    figure_key="auto_eco_logic",
                ),
                Section(
                    "5.7 稳定性与显示连续性测试",
                    [
                        "稳定性测试包括连续停留在首页、连续停留在空气详情页、多次切换页面、多次点击控制按钮和小程序连续刷新。观察重点是屏幕是否闪回初始状态、串口是否出现大量未知事件、ESP32 是否重启、HTTP API 是否断开。",
                        "当前屏幕显示逻辑已经围绕缓存和当前页定向刷新设计。HMI 控件状态也需要配合设置为保持显示，避免屏幕自身在没有新指令时把控件恢复为初始状态。固件和 HMI 文件两侧配合后，才能达到“新数据覆盖旧数据”的效果。",
                    ],
                ),
                Section(
                    "5.8 测试结论",
                    [
                        "测试结果说明，作品已经具备完整的环境感知、现场显示、手机控制、语音控制和自动联动能力。三个交互入口能够围绕同一套设备状态工作，执行器动作可以被直接观察，自动和节能模式也有明确行为。",
                        "当前版本适合完成室内小空间环境管理的核心展示。后续若进入长期使用，可以继续补充云端数据库、更多传感器、真实继电器或电机驱动模块，并增加设备权限和网络安全配置。",
                    ],
                ),
            ],
        ),
        Chapter(
            "第 6 章 应用前景",
            [
                Section(
                    "6.1 家庭与宿舍应用",
                    [
                        "在家庭卧室和客厅中，本作品可以作为桌面环境终端。用户可以直接看串口屏，也可以用手机查看状态。如果空气评分偏低，系统给出通风或净化建议；如果湿度偏低，系统提示加湿。对于夜间使用，节能模式可以一键关闭执行器，减少不必要的运行。",
                        "在宿舍中，空气和湿度变化往往不是某一个人的问题，而是多人共用空间的问题。把设备放在公共位置后，屏幕上的数据可以让大家直观看到当前状态，小程序和语音入口则让控制方式更灵活。作品可以帮助宿舍形成更及时的通风和环境管理习惯。",
                    ],
                ),
                Section(
                    "6.2 教室与办公室应用",
                    [
                        "教室和办公室的共同特点是人员密集、空气变化快。系统可以放在讲台、门口或办公区，作为空气状态提醒设备。管理者看到空气评分下降后，可以开窗、开新风或提醒人员休息。小程序可以给负责维护的人查看近期变化，避免只靠主观感受判断。",
                        "如果后续接入真实风机、净化器或加湿器，系统就能从提醒型设备扩展为控制型设备。当前使用 LED 和舵机是为了清楚展示控制逻辑，真实应用中可以把这些输出换成继电器、MOS 管驱动或标准智能家居接口。",
                    ],
                ),
                Section(
                    "6.3 后续扩展方向",
                    [
                        "传感器方面，后续可以加入 SHT30、PM2.5、CO2、TVOC 等模块，让空气状态更细分。数据方面，可以把最近 30 条记录扩展到云端数据库，形成按小时、按天查看的曲线。交互方面，可以继续完善小程序界面，增加设备命名、家庭房间分组和异常提醒。",
                        "AI 方面，当前重点是语音控制工具调用，后续可以在保证安全的前提下加入更自然的建议，例如根据历史数据提醒“最近晚上湿度偏低，可以提前开低档加湿”。但控制权仍应保持清楚，用户可以随时通过屏幕或小程序关闭自动和节能模式。",
                    ],
                    figure_key="application_scenarios",
                ),
                Section(
                    "6.4 总结",
                    [
                        "本作品完成了基于 ESP32S3 的室内环境智能调节系统。系统可以采集温湿度和空气质量数据，在串口屏上持续显示，也可以通过小程序和 AI 语音控制净化、新风、加湿、自动和节能模式。执行器动作可观察，控制状态统一，适合现场演示和小型空间应用。",
                        "后续扩展可以围绕传感器精度、云端存储、真实设备驱动和更完善的小程序体验展开。整体方向是让室内环境从“凭感觉处理”变成“看得见、控得住、能自动调整”的日常工具。",
                    ],
                ),
            ],
        ),
    ]
    tables = [gpio_table(), api_table(), mcp_table(), test_table()]
    return abstract, chapters, tables


def register_pdf_fonts() -> None:
    pdfmetrics.registerFont(TTFont("DocSong", str(FONT_BODY)))
    pdfmetrics.registerFont(TTFont("DocHei", str(FONT_BOLD)))


class TocDocTemplate(SimpleDocTemplate):
    def afterFlowable(self, flowable):  # noqa: N802 - reportlab hook
        level = getattr(flowable, "_toc_level", None)
        if level is not None:
            self.notify("TOCEntry", (level, flowable.getPlainText(), self.page))


def make_pdf_styles() -> dict[str, ParagraphStyle]:
    base = getSampleStyleSheet()
    return {
        "cover_title": ParagraphStyle("cover_title", parent=base["Title"], fontName="DocHei", fontSize=26, leading=34, alignment=TA_CENTER, spaceAfter=16),
        "cover_sub": ParagraphStyle("cover_sub", parent=base["Normal"], fontName="DocSong", fontSize=16, leading=24, alignment=TA_CENTER, spaceAfter=22),
        "h1": ParagraphStyle("h1", parent=base["Heading1"], fontName="DocHei", fontSize=22, leading=30, spaceBefore=12, spaceAfter=18, wordWrap="CJK"),
        "h2": ParagraphStyle("h2", parent=base["Heading2"], fontName="DocHei", fontSize=16, leading=24, spaceBefore=10, spaceAfter=8, wordWrap="CJK"),
        "body": ParagraphStyle("body", parent=base["BodyText"], fontName="DocSong", fontSize=10.5, leading=19, firstLineIndent=21, alignment=TA_LEFT, spaceAfter=6, wordWrap="CJK"),
        "caption": ParagraphStyle("caption", parent=base["Normal"], fontName="DocSong", fontSize=10, leading=14, alignment=TA_CENTER, spaceBefore=4, spaceAfter=10, wordWrap="CJK"),
        "toc_title": ParagraphStyle("toc_title", parent=base["Heading1"], fontName="DocHei", fontSize=22, leading=30, alignment=TA_CENTER, spaceAfter=22),
        "toc": ParagraphStyle("toc", parent=base["Normal"], fontName="DocSong", fontSize=11, leading=16, wordWrap="CJK"),
        "table": ParagraphStyle("table", parent=base["Normal"], fontName="DocSong", fontSize=9.5, leading=14, wordWrap="CJK"),
        "table_head": ParagraphStyle("table_head", parent=base["Normal"], fontName="DocHei", fontSize=9.5, leading=14, alignment=TA_CENTER, wordWrap="CJK"),
    }


def pdf_paragraph(text: str, style: ParagraphStyle) -> Paragraph:
    return Paragraph(text.replace("\n", "<br/>"), style)


def heading(text: str, style: ParagraphStyle, level: int) -> Paragraph:
    p = Paragraph(text, style)
    p._toc_level = level  # type: ignore[attr-defined]
    return p


def pdf_image(path: Path, max_width: float = 15.5 * cm, max_height: float = 9 * cm) -> RLImage:
    img = Image.open(path)
    w, h = img.size
    scale = min(max_width / w, max_height / h)
    return RLImage(str(path), width=w * scale, height=h * scale)


def pdf_table(table: TableData, styles: dict[str, ParagraphStyle]) -> list:
    data = [[pdf_paragraph(h, styles["table_head"]) for h in table.headers]]
    for row in table.rows:
        data.append([pdf_paragraph(cell, styles["table"]) for cell in row])
    widths = [15.5 * cm / len(table.headers)] * len(table.headers)
    t = Table(data, colWidths=widths, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E9EFEA")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor(GREEN)),
        ("GRID", (0, 0), (-1, -1), 0.6, colors.HexColor("#9AA8A0")),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    return [pdf_paragraph(table.title, styles["caption"]), t, Spacer(1, 8)]


def page_footer(canvas, doc) -> None:
    canvas.saveState()
    canvas.setFont("DocSong", 9)
    canvas.setFillColor(colors.HexColor("#555555"))
    canvas.drawCentredString(A4[0] / 2, 1.05 * cm, str(canvas.getPageNumber()))
    canvas.restoreState()


def generate_pdf(figures: dict[str, Figure], abstract: list[str], chapters: list[Chapter]) -> None:
    register_pdf_fonts()
    styles = make_pdf_styles()
    doc = TocDocTemplate(str(REPORT_PDF), pagesize=A4, rightMargin=2.0 * cm, leftMargin=2.5 * cm, topMargin=2.5 * cm, bottomMargin=2.0 * cm)
    story: list = []

    story.append(Spacer(1, 1.0 * cm))
    story.append(pdf_paragraph("中国大学生计算机设计大赛", styles["cover_title"]))
    story.append(pdf_paragraph("物联网应用类作品技术文档", styles["cover_sub"]))
    story.append(Spacer(1, 1.4 * cm))
    cover_rows = [
        ["作品编号：", ""],
        ["作品名称：", ""],
        ["作者：", ""],
        ["版本编号：", ""],
        ["填写日期：", ""],
    ]
    cover_table = Table(cover_rows, colWidths=[3.2 * cm, 11 * cm], rowHeights=[1.1 * cm] * 5)
    cover_table.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), "DocSong"),
        ("FONTSIZE", (0, 0), (-1, -1), 15),
        ("LINEBELOW", (1, 0), (1, -1), 0.8, colors.black),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    story.append(cover_table)
    story.append(Spacer(1, 1.2 * cm))
    explain = [
        ["填写说明："],
        ["1. 本文档适用于物联网应用类作品。"],
        ["2. 正文采用五号宋体风格，一级标题采用二号黑体风格。"],
        ["3. 文档内容围绕作品实际完成情况填写，重点说明功能、实现和应用价值。"],
        ["4. 封面信息由参赛团队按提交要求自行补充。"],
    ]
    ex_table = Table(explain, colWidths=[14.2 * cm])
    ex_table.setStyle(TableStyle([
        ("BOX", (0, 0), (-1, -1), 0.8, colors.black),
        ("FONTNAME", (0, 0), (-1, -1), "DocSong"),
        ("FONTSIZE", (0, 0), (-1, -1), 10.5),
        ("LEADING", (0, 0), (-1, -1), 16),
        ("LEFTPADDING", (0, 0), (-1, -1), 14),
        ("RIGHTPADDING", (0, 0), (-1, -1), 14),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
    ]))
    story.append(ex_table)
    story.append(PageBreak())

    story.append(pdf_paragraph("喵伴空气管家室内环境智能调节系统", styles["cover_title"]))
    story.append(pdf_paragraph("摘要", styles["h2"]))
    for para in abstract:
        story.append(pdf_paragraph(para, styles["body"]))
    story.append(pdf_paragraph("关键词：ESP32S3；物联网；串口屏；环境监测；智能家居控制；微信小程序；MCP 语音控制", styles["body"]))
    story.append(PageBreak())

    story.append(pdf_paragraph("目录", styles["toc_title"]))
    toc = TableOfContents()
    toc.levelStyles = [
        ParagraphStyle("toc1", fontName="DocSong", fontSize=11, leading=18, leftIndent=0),
        ParagraphStyle("toc2", fontName="DocSong", fontSize=10.5, leading=16, leftIndent=18),
    ]
    story.append(toc)
    story.append(PageBreak())

    for chapter in chapters:
        story.append(heading(chapter.title, styles["h1"], 0))
        for sec in chapter.sections:
            story.append(heading(sec.title, styles["h2"], 1))
            for para in sec.paragraphs:
                story.append(pdf_paragraph(para, styles["body"]))
            if sec.figure_key:
                fig = figures[sec.figure_key]
                story.append(Spacer(1, 6))
                story.append(pdf_image(fig.path))
                story.append(pdf_paragraph(fig.title, styles["caption"]))
            if sec.title == "4.7 工程成果展示":
                for img_name, cap in [
                    ("hardware_overview_1.jpg", "图 11 作品硬件整体图（一）"),
                    ("hardware_overview_2.jpg", "图 12 作品硬件整体图（二）"),
                    ("serial_page_home.png", "图 13 串口屏首页实拍"),
                    ("serial_page_air_score.png", "图 14 串口屏空气详情页实拍"),
                    ("serial_page_ai_settings.png", "图 15 串口屏 AI 与设置页实拍"),
                    ("05_mini_program_flow.png", "图 16 小程序控制界面结构示意"),
                ]:
                    p = ASSET_DIR / img_name
                    if p.exists():
                        story.append(pdf_image(p, max_height=7.5 * cm))
                        story.append(pdf_paragraph(cap, styles["caption"]))
            if sec.table:
                story.extend(pdf_table(sec.table, styles))
            story.append(Spacer(1, 6))
        story.append(PageBreak())

    story.append(heading("附录 A GPIO 与接口表", styles["h1"], 0))
    story.extend(pdf_table(gpio_table(), styles))
    story.append(heading("附录 B HTTP API 与 MCP 工具", styles["h1"], 0))
    story.extend(pdf_table(api_table(), styles))
    story.extend(pdf_table(mcp_table(), styles))
    story.append(heading("附录 C 视频拍摄建议", styles["h1"], 0))
    for para in video_guide_paragraphs():
        story.append(pdf_paragraph(para, styles["body"]))
    story.append(heading("参考文献", styles["h1"], 0))
    refs = [
        "[1] Espressif Systems. ESP-IDF Programming Guide.",
        "[2] Aosong Electronics. DHT11 Temperature and Humidity Sensor Datasheet.",
        "[3] Zhengzhou Winsen Electronics. MQ135 Gas Sensor Technical Manual.",
        "[4] TJC Serial Screen Instruction Set and Development Guide.",
        "[5] 微信公众平台. 小程序开发文档.",
        "[6] Model Context Protocol Specification and FastMCP Usage Documentation.",
    ]
    for ref in refs:
        story.append(pdf_paragraph(ref, styles["body"]))

    doc.multiBuild(story, onFirstPage=page_footer, onLaterPages=page_footer)


def set_docx_font(run, name: str = "宋体", size: int | None = None, bold: bool | None = None) -> None:
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    if size:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold


def set_style_font(style, name: str, size: int, bold: bool = False) -> None:
    style.font.name = name
    style._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    style.font.size = Pt(size)
    style.font.bold = bold


def docx_add_paragraph(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.first_line_indent = Pt(21)
    p.paragraph_format.line_spacing = 1.5
    p.paragraph_format.space_after = Pt(4)
    run = p.add_run(text)
    set_docx_font(run, "宋体", 10)


def docx_add_heading(doc: Document, text: str, level: int) -> None:
    p = doc.add_paragraph(style=f"Heading {level}")
    p.paragraph_format.space_before = Pt(8)
    p.paragraph_format.space_after = Pt(8)
    if level == 1:
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(text)
    set_docx_font(run, "黑体", 22 if level == 1 else 15, True)


def docx_add_picture(doc: Document, path: Path, caption: str, width_in: float = 5.7) -> None:
    if not path.exists():
        return
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(path), width=Inches(width_in))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cap.add_run(caption)
    set_docx_font(run, "宋体", 10)


def docx_add_table(doc: Document, table: TableData) -> None:
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cap.add_run(table.title)
    set_docx_font(run, "宋体", 10)
    t = doc.add_table(rows=1, cols=len(table.headers))
    t.style = "Table Grid"
    for i, h in enumerate(table.headers):
        cell = t.rows[0].cells[i]
        cell.text = h
        for p in cell.paragraphs:
            for r in p.runs:
                set_docx_font(r, "黑体", 9, True)
    for row in table.rows:
        cells = t.add_row().cells
        for i, cell_text in enumerate(row):
            cells[i].text = cell_text
            for p in cells[i].paragraphs:
                for r in p.runs:
                    set_docx_font(r, "宋体", 9)


def generate_docx(figures: dict[str, Figure], abstract: list[str], chapters: list[Chapter]) -> None:
    doc = Document()
    for section in doc.sections:
        section.page_width = Cm(21)
        section.page_height = Cm(29.7)
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.0)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.0)
    set_style_font(doc.styles["Normal"], "宋体", 10)
    set_style_font(doc.styles["Heading 1"], "黑体", 22, True)
    set_style_font(doc.styles["Heading 2"], "黑体", 15, True)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("中国大学生计算机设计大赛")
    set_docx_font(run, "黑体", 22, True)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("物联网应用类作品技术文档")
    set_docx_font(run, "宋体", 18)
    doc.add_paragraph()
    for label in ["作品编号：", "作品名称：", "作者：", "版本编号：", "填写日期："]:
        p = doc.add_paragraph()
        run = p.add_run(label + "____________________________")
        set_docx_font(run, "宋体", 15)
    doc.add_page_break()

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("喵伴空气管家室内环境智能调节系统")
    set_docx_font(run, "黑体", 18, True)
    docx_add_heading(doc, "摘要", 2)
    for para in abstract:
        docx_add_paragraph(doc, para)
    docx_add_paragraph(doc, "关键词：ESP32S3；物联网；串口屏；环境监测；智能家居控制；微信小程序；MCP 语音控制")
    doc.add_page_break()

    docx_add_heading(doc, "目录", 1)
    for chapter in chapters:
        docx_add_paragraph(doc, chapter.title)
        for sec in chapter.sections:
            docx_add_paragraph(doc, "    " + sec.title)
    doc.add_page_break()

    for chapter in chapters:
        docx_add_heading(doc, chapter.title, 1)
        for sec in chapter.sections:
            docx_add_heading(doc, sec.title, 2)
            for para in sec.paragraphs:
                docx_add_paragraph(doc, para)
            if sec.figure_key:
                fig = figures[sec.figure_key]
                docx_add_picture(doc, fig.path, fig.title)
            if sec.title == "4.7 工程成果展示":
                for img_name, cap in [
                    ("hardware_overview_1.jpg", "图 11 作品硬件整体图（一）"),
                    ("hardware_overview_2.jpg", "图 12 作品硬件整体图（二）"),
                    ("serial_page_home.png", "图 13 串口屏首页实拍"),
                    ("serial_page_air_score.png", "图 14 串口屏空气详情页实拍"),
                    ("serial_page_ai_settings.png", "图 15 串口屏 AI 与设置页实拍"),
                    ("05_mini_program_flow.png", "图 16 小程序控制界面结构示意"),
                ]:
                    docx_add_picture(doc, ASSET_DIR / img_name, cap, width_in=5.4)
            if sec.table:
                docx_add_table(doc, sec.table)
        doc.add_page_break()

    docx_add_heading(doc, "附录 A GPIO 与接口表", 1)
    docx_add_table(doc, gpio_table())
    docx_add_heading(doc, "附录 B HTTP API 与 MCP 工具", 1)
    docx_add_table(doc, api_table())
    docx_add_table(doc, mcp_table())
    docx_add_heading(doc, "附录 C 视频拍摄建议", 1)
    for para in video_guide_paragraphs():
        docx_add_paragraph(doc, para)
    docx_add_heading(doc, "参考文献", 1)
    for ref in [
        "[1] Espressif Systems. ESP-IDF Programming Guide.",
        "[2] Aosong Electronics. DHT11 Temperature and Humidity Sensor Datasheet.",
        "[3] Zhengzhou Winsen Electronics. MQ135 Gas Sensor Technical Manual.",
        "[4] TJC Serial Screen Instruction Set and Development Guide.",
        "[5] 微信公众平台. 小程序开发文档.",
        "[6] Model Context Protocol Specification and FastMCP Usage Documentation.",
    ]:
        docx_add_paragraph(doc, ref)
    doc.save(REPORT_DOCX)


def ppt_rgb(hex_color: str) -> PptRGB:
    hex_color = hex_color.strip("#")
    return PptRGB(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_ppt_title(slide, text: str, dark: bool = False) -> None:
    box = slide.shapes.add_textbox(PptInches(0.55), PptInches(0.35), PptInches(12.1), PptInches(0.65))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = PptPt(34)
    p.font.bold = True
    p.font.color.rgb = ppt_rgb("#FFFFFF" if dark else GREEN)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 16, bold: bool = False, color: str = DARK) -> None:
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = PptPt(size)
    p.font.bold = bold
    p.font.color.rgb = ppt_rgb(color)


def add_card(slide, x: float, y: float, w: float, h: float, title_text: str, body: str, accent: str = ORANGE) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb("#FFFFFF")
    shape.line.color.rgb = ppt_rgb("#D8DED8")
    add_text(slide, title_text, x + 0.25, y + 0.18, w - 0.5, 0.35, 17, True, GREEN)
    add_text(slide, body, x + 0.25, y + 0.72, w - 0.5, h - 0.85, 13, False, DARK)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptInches(x), PptInches(y), PptInches(0.07), PptInches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ppt_rgb(accent)
    bar.line.fill.background()


def set_slide_bg(slide, color: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ppt_rgb(color)


def add_image(slide, path: Path, x: float, y: float, w: float | None = None, h: float | None = None) -> None:
    if not path.exists():
        return
    kwargs = {}
    if w is not None:
        kwargs["width"] = PptInches(w)
    if h is not None:
        kwargs["height"] = PptInches(h)
    slide.shapes.add_picture(str(path), PptInches(x), PptInches(y), **kwargs)


def generate_pptx(figures: dict[str, Figure]) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, GREEN)
    add_text(slide, "喵伴空气管家", 0.8, 1.25, 7.4, 0.8, 42, True, "#FFFFFF")
    add_text(slide, "室内环境智能调节系统", 0.82, 2.05, 6.4, 0.45, 21, False, "#FFFFFF")
    add_text(slide, "温湿度和空气质量检测 · 串口屏现场交互 · 小程序控制 · AI 语音控制", 0.82, 5.9, 9.0, 0.4, 18, False, "#FFFFFF")
    add_image(slide, ASSET_DIR / "hardware_overview_1.jpg", 8.3, 1.2, w=4.2)

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT)
    add_ppt_title(slide, "1. 作品解决的问题")
    add_card(slide, 0.75, 1.45, 3.55, 4.6, "环境变化不直观", "很多室内空间只能靠感觉判断闷、干、热，发现问题时往往已经晚了。")
    add_card(slide, 4.85, 1.45, 3.55, 4.6, "控制入口分散", "屏幕、手机、语音如果各管各的，状态容易不一致，用户也不好理解。")
    add_card(slide, 8.95, 1.45, 3.55, 4.6, "手动操作容易忘", "需要一个自动模式处理明显异常，也需要一个节能模式快速关闭设备。")

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT)
    add_ppt_title(slide, "2. 我们做了什么")
    add_image(slide, figures["system_architecture"].path, 0.75, 1.25, w=7.2)
    add_card(slide, 8.45, 1.35, 3.9, 1.35, "感知", "DHT11 读取温湿度，MQ135 读取空气质量模拟量。")
    add_card(slide, 8.45, 2.95, 3.9, 1.35, "交互", "串口屏、小程序、AI 语音都接入同一套状态。")
    add_card(slide, 8.45, 4.55, 3.9, 1.35, "执行", "红灯表示净化，蓝灯表示加湿，舵机扇叶表示新风。")

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT)
    add_ppt_title(slide, "3. 硬件和接口")
    add_image(slide, figures["hardware_wiring"].path, 0.55, 1.15, w=7.4)
    add_text(slide, "关键引脚", 8.35, 1.35, 3.7, 0.4, 22, True, GREEN)
    add_text(slide, "GPIO18: DHT11\nGPIO1: MQ135 ADC\nGPIO41/42: TJC 串口屏\nGPIO13: 净化红灯\nGPIO14: 加湿蓝灯\nGPIO21: 新风舵机扇叶", 8.35, 2.0, 4.2, 3.2, 17, False, DARK)

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT)
    add_ppt_title(slide, "4. 串口屏现场交互")
    add_image(slide, ASSET_DIR / "serial_page_home.png", 0.75, 1.25, w=3.55)
    add_image(slide, ASSET_DIR / "serial_page_air_score.png", 4.9, 1.25, w=3.55)
    add_image(slide, ASSET_DIR / "serial_page_ai_settings.png", 9.05, 1.25, w=3.55)
    add_text(slide, "首页看状态，详情页看曲线，AI/设置页切换手动环境。数据保持在屏幕上，新数据来时覆盖更新。", 0.85, 6.05, 11.8, 0.55, 18, False, DARK)

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT)
    add_ppt_title(slide, "5. 小程序控制")
    add_image(slide, figures["mini_program_flow"].path, 0.75, 1.15, w=7.2)
    add_card(slide, 8.35, 1.35, 3.95, 1.25, "读取状态", "GET /api/state 获取温湿度、空气评分和设备档位。")
    add_card(slide, 8.35, 2.85, 3.95, 1.25, "控制设备", "POST /api/device 切换净化、新风、加湿 0-3 档。")
    add_card(slide, 8.35, 4.35, 3.95, 1.25, "查看历史", "GET /api/history 显示最近 30 条空气评分。")

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT)
    add_ppt_title(slide, "6. AI 语音控制")
    add_image(slide, figures["mcp_bridge_flow"].path, 0.75, 1.25, w=7.4)
    add_text(slide, "语音不是直接操作 GPIO，而是调用 MCP 工具，再由本地桥接程序转发到 ESP32 HTTP API。这样屏幕、小程序和语音控制看到的是同一套状态。", 8.45, 1.65, 3.85, 2.4, 17, False, DARK)
    add_text(slide, "示例：打开净化器二档\n示例：模拟污染环境\n示例：进入节能模式", 8.45, 4.55, 3.85, 1.2, 18, True, ORANGE)

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT)
    add_ppt_title(slide, "7. 自动与节能模式")
    add_image(slide, figures["auto_eco_logic"].path, 0.75, 1.15, w=7.4)
    add_card(slide, 8.45, 1.45, 3.8, 1.25, "自动模式", "根据湿度、空气评分和温度调整执行器档位。")
    add_card(slide, 8.45, 3.05, 3.8, 1.25, "节能模式", "关闭自动模式，并把三个执行器档位置 0。")
    add_card(slide, 8.45, 4.65, 3.8, 1.25, "现场展示", "用舒适、高温、干燥、污染场景验证联动。")

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT)
    add_ppt_title(slide, "8. 测试与完成情况")
    add_image(slide, figures["test_workflow"].path, 0.65, 1.05, w=7.2)
    add_text(slide, "已完成验证", 8.3, 1.25, 3.5, 0.4, 22, True, GREEN)
    add_text(slide, "传感器读数\n串口屏显示和触摸\n小程序状态读取和控制\nAI/MCP 工具调用\n自动/节能联动\n近期空气质量记录", 8.3, 1.95, 4.1, 3.2, 17, False, DARK)

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, GREEN)
    add_ppt_title(slide, "9. 应用价值与扩展", dark=True)
    add_image(slide, figures["application_scenarios"].path, 0.65, 1.2, w=7.25)
    add_text(slide, "落地价值", 8.35, 1.4, 3.8, 0.4, 23, True, "#FFFFFF")
    add_text(slide, "让室内环境从“凭感觉处理”变成“看得见、控得住、能自动调整”。\n\n后续可以接入云端记录、PM2.5/CO2/TVOC 传感器和真实净化器、风机、加湿器控制。", 8.35, 2.05, 4.1, 2.6, 17, False, "#FFFFFF")
    add_text(slide, "适用：家庭、宿舍、教室、小办公室", 8.35, 5.65, 4.2, 0.4, 19, True, ORANGE)

    prs.save(PPTX_PATH)


def video_guide_paragraphs() -> list[str]:
    return [
        "视频建议控制在 3 分钟以内，重点放在作品演示。开头 10 到 15 秒展示实物整体和一句话定位：这是一套能检测室内温湿度和空气质量，并支持屏幕、小程序、AI 语音控制的环境调节系统。",
        "第一个演示段落展示串口屏。先拍首页，说明温度、湿度、空气评分和建议会持续显示；再切到空气详情页，展示空气评分和近期曲线；最后切到智能家居或 AI 设置相关页面，说明现场可以直接触摸操作。",
        "第二个演示段落展示执行器。依次点击净化、新风、加湿，画面要拍到红色 LED、舵机扇叶和蓝色 LED 的变化。讲述时不要只说按钮亮了，要说它们分别对应净化、新风和加湿，并且档位可以通过亮度或舵机摆动表现出来。",
        "第三个演示段落展示小程序。镜头拍小程序连接 ESP32 地址、读取状态、切换设备档位、查看近期空气质量记录。这里要强调手机端和串口屏控制的是同一套 ESP32 状态。",
        "第四个演示段落展示 AI 语音控制。选择三条最稳的口令：打开净化器二档、开启自动模式、模拟污染环境。每条口令后镜头要回到实物或小程序状态，让观众看到控制结果。",
        "最后 10 到 15 秒收束应用价值。可以说系统适合家庭、宿舍、教室和小办公室，后续可以接入真实净化器、加湿器、风机和更多空气传感器，让小空间环境管理更直观。",
    ]


def write_video_guide() -> None:
    content = ["# 视频拍摄与答辩讲述建议", ""]
    for i, para in enumerate(video_guide_paragraphs(), start=1):
        content.append(f"{i}. {para}")
        content.append("")
    content.append("## 3 分钟视频推荐节奏")
    content.extend([
        "",
        "- 0:00-0:15 实物整体和作品定位。",
        "- 0:15-0:45 串口屏首页、空气详情和曲线。",
        "- 0:45-1:20 净化、新风、加湿、自动、节能控制。",
        "- 1:20-2:00 小程序读取状态、控制档位、查看历史。",
        "- 2:00-2:40 AI 语音控制三条口令。",
        "- 2:40-3:00 应用场景和可扩展方向。",
    ])
    VIDEO_GUIDE.write_text("\n".join(content), encoding="utf-8")


def main() -> None:
    ensure_dirs()
    figures = create_figures()
    abstract, chapters, _ = build_chapters()
    generate_docx(figures, abstract, chapters)
    generate_pdf(figures, abstract, chapters)
    generate_pptx(figures)
    write_video_guide()
    print("generated")
    print(REPORT_DOCX)
    print(REPORT_PDF)
    print(PPTX_PATH)
    print(VIDEO_GUIDE)


if __name__ == "__main__":
    main()
