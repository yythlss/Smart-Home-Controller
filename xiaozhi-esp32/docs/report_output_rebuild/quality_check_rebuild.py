from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


ROOT = Path(__file__).resolve().parent
FINAL = ROOT / "final"
ASSETS = ROOT / "assets"


def find_one(pattern: str) -> Path:
    matches = list(FINAL.glob(pattern))
    if len(matches) != 1:
        raise AssertionError(f"expected one {pattern}, found {len(matches)}")
    return matches[0]


def docx_text(path: Path) -> str:
    doc = Document(str(path))
    parts: list[str] = []
    for p in doc.paragraphs:
        parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)


def pptx_text(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)


def pdf_text(reader: PdfReader) -> str:
    parts: list[str] = []
    for page in reader.pages[:8]:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            parts.append("")
    return "\n".join(parts)


def assert_contains(text: str, needle: str) -> None:
    if needle not in text:
        raise AssertionError(f"missing required text: {needle}")


def assert_absent(text: str, phrase: str) -> None:
    if phrase in text:
        raise AssertionError(f"banned phrase found: {phrase}")


def main() -> None:
    docx = find_one("*.docx")
    pdf = find_one("*.pdf")
    pptx = find_one("*.pptx")
    guide = find_one("*.md")

    dtext = docx_text(docx)
    ptext = pptx_text(pptx)
    reader = PdfReader(str(pdf))
    extracted_pdf = pdf_text(reader)
    combined = "\n".join([dtext, ptext, extracted_pdf])

    if len(reader.pages) < 20:
        raise AssertionError(f"PDF page count too low: {len(reader.pages)}")

    slide_count = len(Presentation(str(pptx)).slides)
    if not 8 <= slide_count <= 10:
        raise AssertionError(f"PPT slide count out of range: {slide_count}")

    for phrase in ["仅用于竞赛演示", "演示级", "原型演示级", "最难的是", "TBD", "TODO", "lorem", "ipsum"]:
        assert_absent(combined, phrase)

    for required in [
        "第 1 章 作品概述",
        "第 2 章 需求分析",
        "第 3 章 技术方案",
        "第 4 章 方案实现",
        "第 5 章 测试报告",
        "第 6 章 应用前景",
        "串口屏",
        "小程序",
        "AI/MCP",
        "自动模式",
        "节能模式",
        "GPIO13",
        "GPIO14",
        "GPIO21",
        "DHT11",
        "MQ135",
        "/api/state",
        "/api/device",
    ]:
        assert_contains(dtext, required)

    custom_figures = list(ASSETS.glob("[0-9][0-9]_*.png"))
    if len(custom_figures) < 9:
        raise AssertionError(f"custom figure count too low: {len(custom_figures)}")

    if guide.stat().st_size < 1000:
        raise AssertionError("video guide is unexpectedly short")

    print("ALL CHECKS PASSED")
    print(f"docx={docx}")
    print(f"pdf_pages={len(reader.pages)}")
    print(f"ppt_slides={slide_count}")
    print(f"custom_figures={len(custom_figures)}")


if __name__ == "__main__":
    main()
