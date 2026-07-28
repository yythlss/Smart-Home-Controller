from pathlib import Path
import textwrap

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import (
    Image as PdfImage,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "docs" / "report_output"
ASSET_DIR = OUT_DIR / "assets"
FIG_DIR = OUT_DIR / "figures_final"
FINAL_DIR = OUT_DIR / "final"
FIG_DIR.mkdir(parents=True, exist_ok=True)
FINAL_DIR.mkdir(parents=True, exist_ok=True)

REPORT_DOCX = FINAL_DIR / "喵伴空气管家作品设计报告_2026-07-08.docx"
REPORT_PDF = FINAL_DIR / "喵伴空气管家作品设计报告_2026-07-08.pdf"
PPTX_OUT = FINAL_DIR / "喵伴空气管家答辩PPT_2026-07-08.pptx"
GUIDE_MD = FINAL_DIR / "视频讲稿与答辩建议_2026-07-08.md"
GUIDE_DOCX = FINAL_DIR / "视频讲稿与答辩建议_2026-07-08.docx"

IMAGES = {
    "home": ASSET_DIR / "serial_page_home.png",
    "air": ASSET_DIR / "serial_page_air_score.png",
    "ai": ASSET_DIR / "serial_page_ai_settings.png",
    "hw1": ASSET_DIR / "hardware_overview_1.jpg",
    "hw2": ASSET_DIR / "hardware_overview_2.jpg",
}


def pil_font(size=28, bold=False):
    candidates = [
        Path("C:/Windows/Fonts/simhei.ttf") if bold else Path("C:/Windows/Fonts/simsun.ttc"),
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size)
    return ImageFont.load_default()


def wrap_cn(text, width):
    lines, cur = [], ""
    for ch in text:
        cur += ch
        if len(cur) >= width:
            lines.append(cur)
            cur = ""
    if cur:
        lines.append(cur)
    return lines


def draw_box(draw, xy, title, lines=None, fill="#F7FAFC", outline="#2F5597"):
    x1, y1, x2, y2 = xy
    draw.rounded_rectangle(xy, radius=18, fill=fill, outline=outline, width=3)
    draw.text((x1 + 18, y1 + 14), title, fill="#17365D", font=pil_font(28, True))
    y = y1 + 58
    for line in lines or []:
        for part in wrap_cn(line, 16):
            draw.text((x1 + 18, y), part, fill="#222222", font=pil_font(21))
            y += 30


def draw_arrow(draw, start, end, color="#2F5597"):
    import math

    draw.line([start, end], fill=color, width=4)
    x1, y1 = start
    x2, y2 = end
    angle = math.atan2(y2 - y1, x2 - x1)
    pts = [
        end,
        (x2 + 16 * math.cos(angle + 2.55), y2 + 16 * math.sin(angle + 2.55)),
        (x2 + 16 * math.cos(angle - 2.55), y2 + 16 * math.sin(angle - 2.55)),
    ]
    draw.polygon(pts, fill=color)


def make_figures():
    arch = FIG_DIR / "fig1_system_architecture.png"
    hw = FIG_DIR / "fig2_hardware_connection.png"
    sw = FIG_DIR / "fig3_software_flow.png"
    interaction = FIG_DIR / "fig4_interaction_channels.png"

    im = Image.new("RGB", (1500, 820), "white")
    d = ImageDraw.Draw(im)
    d.text((50, 30), "系统总体架构图", fill="#17365D", font=pil_font(38, True))
    draw_box(d, (70, 150, 370, 330), "环境感知层", ["DHT11 温湿度", "MQ135 空气质量原始值", "5 秒周期采样"])
    draw_box(d, (70, 460, 370, 640), "人机交互层", ["串口屏触摸", "微信小程序", "AI 语音 MCP"])
    draw_box(d, (560, 250, 940, 530), "ESP32S3 控制核心", ["数据校验与状态保持", "自动/节能策略", "HTTP API 与 MCP 工具", "PWM/舵机输出"])
    draw_box(d, (1130, 130, 1430, 330), "执行反馈层", ["红 LED：净化档位", "蓝 LED：加湿档位", "舵机扇叶：新风档位"])
    draw_box(d, (1130, 490, 1430, 690), "显示反馈层", ["温湿度与空气等级", "空气质量曲线", "运行建议与模式状态"])
    for a, b in [((370, 240), (560, 340)), ((370, 550), (560, 440)), ((940, 350), (1130, 230)), ((940, 450), (1130, 590)), ((1130, 590), (940, 475)), ((560, 460), (370, 555))]:
        draw_arrow(d, a, b)
    im.save(arch)

    im = Image.new("RGB", (1500, 820), "white")
    d = ImageDraw.Draw(im)
    d.text((50, 30), "硬件连接关系图", fill="#17365D", font=pil_font(38, True))
    draw_box(d, (560, 270, 940, 520), "ESP32S3 开发板", ["UART2 / ADC / GPIO / LEDC", "运行小智框架与控制固件"])
    for xy, title, lines in [
        ((80, 100, 400, 250), "TJC 串口屏", ["GPIO41 TX -> 屏 RX", "GPIO42 RX <- 屏 TX", "9600 8N1"]),
        ((80, 330, 400, 480), "DHT11", ["GPIO18 DATA", "温湿度采集"]),
        ((80, 560, 400, 710), "MQ135", ["GPIO1 / ADC1_CH0", "空气质量原始值"]),
        ((1100, 100, 1420, 250), "红色 LED", ["GPIO13 PWM", "净化档位指示"]),
        ((1100, 330, 1420, 480), "蓝色 LED", ["GPIO14 PWM", "加湿档位指示"]),
        ((1100, 560, 1420, 710), "SG90 舵机扇叶", ["GPIO21 50Hz PWM", "0-180° 往复摆动"]),
    ]:
        draw_box(d, xy, title, lines)
    for p in [(400, 175), (400, 405), (400, 635)]:
        draw_arrow(d, p, (560, 395))
    for p in [(1100, 175), (1100, 405), (1100, 635)]:
        draw_arrow(d, (940, 395), p)
    im.save(hw)

    im = Image.new("RGB", (1500, 820), "white")
    d = ImageDraw.Draw(im)
    d.text((50, 30), "软件数据流与控制流程图", fill="#17365D", font=pil_font(38, True))
    boxes = [
        ((70, 150, 340, 300), "传感器采样", ["DHT11 / MQ135", "异常数据过滤"]),
        ((430, 150, 700, 300), "状态融合", ["空气评分", "舒适度与建议", "最近有效值保持"]),
        ((790, 150, 1060, 300), "串口屏刷新", ["只更新控件值", "曲线追加点", "减少闪烁"]),
        ((1150, 150, 1430, 300), "屏幕显示", ["首页", "空气详情", "智能家居", "AI 设置"]),
        ((430, 470, 700, 620), "统一控制器", ["净化/新风/加湿", "自动/节能互斥"]),
        ((790, 470, 1060, 620), "执行器输出", ["LED PWM 亮度", "舵机角度档位"]),
        ((70, 470, 340, 620), "控制入口", ["屏幕按键", "小程序 HTTP", "AI MCP 工具"]),
    ]
    for xy, title, lines in boxes:
        draw_box(d, xy, title, lines)
    for a, b in [((340, 225), (430, 225)), ((700, 225), (790, 225)), ((1060, 225), (1150, 225)), ((925, 300), (925, 470)), ((790, 545), (700, 545)), ((340, 545), (430, 545)), ((700, 545), (790, 545)), ((565, 470), (565, 300))]:
        draw_arrow(d, a, b)
    im.save(sw)

    im = Image.new("RGB", (1500, 760), "white")
    d = ImageDraw.Draw(im)
    d.text((50, 30), "三种交互方式服务同一套设备状态", fill="#17365D", font=pil_font(38, True))
    draw_box(d, (90, 170, 430, 530), "串口屏", ["现场可见", "触摸切换页面", "直接控制设备"], fill="#EAF4FF")
    draw_box(d, (580, 170, 920, 530), "微信小程序", ["同一局域网访问", "查看历史曲线", "远程调档与模式"], fill="#EAF7EA", outline="#3C7D3E")
    draw_box(d, (1070, 170, 1410, 530), "AI 语音 / MCP", ["语音转工具调用", "桥接到 ESP32 HTTP API", "适合免手操作"], fill="#FFF4E5", outline="#C55A11")
    draw_arrow(d, (430, 350), (580, 350))
    draw_arrow(d, (920, 350), (1070, 350))
    d.text((505, 625), "统一进入 SmartHomeController，保证显示、控制和执行结果一致", fill="#333333", font=pil_font(30, True))
    im.save(interaction)
    return arch, hw, sw, interaction


def set_run_font(run, size=None, bold=None, color=None, name="宋体"):
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    if size:
        run.font.size = Pt(size)
    if bold is not None:
        run.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)


def add_page_number(paragraph):
    run = paragraph.add_run()
    fld_begin = OxmlElement("w:fldChar")
    fld_begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = "PAGE"
    fld_end = OxmlElement("w:fldChar")
    fld_end.set(qn("w:fldCharType"), "end")
    run._r.append(fld_begin)
    run._r.append(instr)
    run._r.append(fld_end)


def style_document(doc):
    sec = doc.sections[0]
    sec.page_width = Cm(21)
    sec.page_height = Cm(29.7)
    sec.top_margin = Cm(2.5)
    sec.bottom_margin = Cm(2.0)
    sec.left_margin = Cm(2.5)
    sec.right_margin = Cm(2.0)
    sec.header_distance = Cm(1.5)
    sec.footer_distance = Cm(1.2)

    normal = doc.styles["Normal"]
    normal.font.name = "宋体"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    normal.font.size = Pt(12)

    for name, size in [("Heading 1", 16), ("Heading 2", 14), ("Heading 3", 12)]:
        st = doc.styles[name]
        st.font.name = "黑体"
        st._element.rPr.rFonts.set(qn("w:eastAsia"), "黑体")
        st.font.size = Pt(size)
        st.font.bold = True
        st.font.color.rgb = RGBColor(0, 0, 0)

    header = sec.header.paragraphs[0]
    header.text = "喵伴空气管家作品设计报告"
    header.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in header.runs:
        set_run_font(run, 9)

    footer = sec.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_page_number(footer)


def add_p(doc, text="", style="Normal", align=None, size=None, bold=None, first_line=True):
    p = doc.add_paragraph(style=style)
    if align is not None:
        p.alignment = align
    if first_line and style == "Normal" and align is None:
        p.paragraph_format.first_line_indent = Pt(24)
    p.paragraph_format.line_spacing = 1.25
    p.paragraph_format.space_after = Pt(5)
    if text:
        r = p.add_run(text)
        set_run_font(r, size=size, bold=bold)
    return p


def h1(doc, text):
    p = add_p(doc, text, "Heading 1", WD_ALIGN_PARAGRAPH.CENTER, 16, True, False)
    p.paragraph_format.space_before = Pt(12)
    p.paragraph_format.space_after = Pt(10)


def h2(doc, text):
    p = add_p(doc, text, "Heading 2", None, 14, True, False)
    p.paragraph_format.space_before = Pt(8)
    p.paragraph_format.space_after = Pt(6)


def h3(doc, text):
    p = add_p(doc, text, "Heading 3", None, 12, True, False)
    p.paragraph_format.space_before = Pt(6)


def caption(doc, text):
    p = add_p(doc, text, "Normal", WD_ALIGN_PARAGRAPH.CENTER, 10, False, False)
    p.paragraph_format.space_after = Pt(8)


def add_doc_image(doc, path, width_cm=14.5):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(path), width=Cm(width_cm))
    p.paragraph_format.space_after = Pt(2)


def shade(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def cell_text(cell, text, bold=False, align=WD_ALIGN_PARAGRAPH.CENTER):
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = align
    p.paragraph_format.line_spacing = 1.15
    r = p.add_run(text)
    set_run_font(r, 10, bold)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def add_table(doc, headers, rows):
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    for i, h in enumerate(headers):
        cell_text(table.rows[0].cells[i], h, True)
        shade(table.rows[0].cells[i], "D9EAF7")
    for row in rows:
        cells = table.add_row().cells
        for i, val in enumerate(row):
            cell_text(cells[i], str(val), False, WD_ALIGN_PARAGRAPH.LEFT if i > 0 else WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()
    return table


def add_toc(doc):
    h1(doc, "目录")
    entries = [
        "摘要",
        "第一部分 作品概述",
        "1.1 功能与特性",
        "1.2 应用领域",
        "1.3 主要技术特点",
        "1.4 主要性能指标",
        "1.5 主要创新点",
        "1.6 设计流程",
        "第二部分 系统组成及功能说明",
        "2.1 整体介绍",
        "2.2 硬件系统介绍",
        "2.3 软件系统介绍",
        "第三部分 完成情况及性能参数",
        "3.1 整体介绍",
        "3.2 工程成果",
        "3.3 特性成果",
        "第四部分 总结",
        "4.1 可扩展之处",
        "4.2 心得体会",
        "第五部分 参考文献",
    ]
    for e in entries:
        add_p(doc, e, first_line=False)
    doc.add_page_break()


def build_report(figs):
    arch, hw, sw, interaction = figs
    doc = Document()
    style_document(doc)

    p = add_p(doc, "喵伴空气管家", "Heading 1", WD_ALIGN_PARAGRAPH.CENTER, 20, True, False)
    p.paragraph_format.space_before = Pt(80)
    add_p(doc, "室内环境监测与智能调节系统作品设计报告", "Normal", WD_ALIGN_PARAGRAPH.CENTER, 16, True, False)
    add_p(doc, "适用于作品申报与现场答辩材料", "Normal", WD_ALIGN_PARAGRAPH.CENTER, 11, False, False)
    add_p(doc, "报告类型：作品报告", "Normal", WD_ALIGN_PARAGRAPH.CENTER, 12, False, False)
    add_p(doc, "作品名称：喵伴空气管家", "Normal", WD_ALIGN_PARAGRAPH.CENTER, 12, False, False)
    doc.add_page_break()

    add_toc(doc)

    h1(doc, "摘 要")
    add_p(doc, "本作品面向宿舍、家庭、办公室等小型室内空间，设计并实现了一套“能监测、能显示、能控制、能语音交互”的室内环境智能调节系统。系统以 ESP32S3 为控制核心，接入 DHT11 温湿度传感器和 MQ135 空气质量传感器，使用 TJC 串口屏持续显示温度、湿度、空气等级、空气评分和近期变化曲线；同时通过红色 LED、蓝色 LED 和 SG90 舵机扇叶分别模拟净化、加湿和新风执行机构。用户可以在串口屏上直接控制设备，也可以通过微信小程序在同一局域网内查看状态、调节档位和查看历史数据，还可以通过 AI 语音/MCP 工具完成免手操作。系统设置自动模式和节能模式：自动模式根据温湿度和空气质量联动开启净化、新风或加湿，节能模式关闭自动控制和执行器，适合离开房间或低功耗展示场景。作品的重点不是堆叠单个模块，而是把环境采集、屏幕显示、手机控制、语音控制和执行反馈串成一个完整闭环，让室内环境变化和控制效果都能被直观看到。")
    add_p(doc, "关键词：ESP32S3；室内空气监测；串口屏；微信小程序；MCP 语音控制；智能家居")

    h1(doc, "第一部分  作品概述")
    h2(doc, "1.1 功能与特性")
    add_p(doc, "作品完成了室内温湿度与空气质量采集、串口屏可视化显示、屏幕触摸控制、小程序局域网控制、AI 语音控制和自动联动调节等功能。串口屏首页展示温度、湿度和空气状态，空气详情页展示空气评分、MQ135 原始值、舒适度建议和近期空气质量曲线，智能家居页提供净化、新风、加湿、自动、节能等按键。执行反馈采用红色 LED 表示净化档位、蓝色 LED 表示加湿档位、舵机扇叶摆动表示新风档位，评委可以直接看到控制结果。")
    add_doc_image(doc, arch, 15)
    caption(doc, "图 1 系统总体架构图")

    h2(doc, "1.2 应用领域")
    add_p(doc, "本作品适用于宿舍、家庭、办公室、教室、实验室等小型室内空间的环境监测与智能联动演示。对于实际应用，它可以提醒用户关注温湿度和空气状态，并在空气较差或环境干燥时自动给出调节建议；对于课程设计和竞赛展示，它把传感器、串口屏、网络接口、语音工具和执行器整合在一起，能够清楚展示嵌入式系统从感知到控制的完整流程。后续更换高精度 PM2.5、CO2、TVOC 传感器并接入真实风机、加湿器或净化器后，可进一步扩展为可用的家庭环境管理终端。")

    h2(doc, "1.3 主要技术特点")
    add_p(doc, "（1）以 ESP32S3 为统一控制核心，完成传感器采集、串口屏通信、HTTP API、MCP 工具注册和执行器控制。")
    add_p(doc, "（2）TJC 串口屏负责静态页面和控件显示，固件只更新动态数据，避免频繁整屏重绘带来的闪烁。")
    add_p(doc, "（3）微信小程序采用局域网 HTTP 方案，可直接访问 ESP32S3 的状态、历史、设备、模式和环境模拟接口。")
    add_p(doc, "（4）AI 语音控制通过 MCP 工具把语音意图转为设备控制命令，最终仍进入同一套 SmartHomeController 状态管理逻辑。")
    add_p(doc, "（5）自动模式、节能模式和手动环境模拟统一接入控制规则，方便答辩时演示污染、干燥、高温等不同场景。")

    h2(doc, "1.4 主要性能指标")
    add_table(doc, ["指标项", "当前实现"], [
        ["主控芯片", "ESP32S3，基于现有小智 ESP32 工程扩展"],
        ["传感器", "DHT11 温湿度；MQ135 模拟空气质量原始值"],
        ["数据刷新", "传感器任务约 5 秒采样一次，屏幕保持上一帧有效数据"],
        ["串口屏通信", "UART2，GPIO41/GPIO42，9600bps，8N1"],
        ["执行器", "GPIO13 红 LED、GPIO14 蓝 LED、GPIO21 SG90 舵机扇叶"],
        ["设备档位", "净化/新风/加湿均支持 0-3 档"],
        ["小程序接口", "GET /api/state、GET /api/history、POST /api/device、POST /api/mode、POST /api/environment"],
        ["历史数据", "最多缓存 30 条近期空气质量样本，用于曲线展示"],
    ])
    caption(doc, "表 1 主要性能指标")

    h2(doc, "1.5 主要创新点")
    add_p(doc, "（1）把串口屏、小程序和 AI 语音三种交互方式统一到同一套设备状态，避免多入口控制结果不一致。")
    add_p(doc, "（2）用 LED 亮度和舵机扇叶动作表达设备档位，比只显示文字更适合现场演示。")
    add_p(doc, "（3）通过手动环境模拟功能快速构造高温、干燥、污染等场景，让自动模式不依赖现场空气变化也能展示。")
    add_p(doc, "（4）串口屏采用控件化更新和曲线控件追加点，使界面更连续，减少数据刷新带来的闪烁感。")

    h2(doc, "1.6 设计流程")
    add_p(doc, "设计流程为：先确认 ESP32S3 可用引脚和硬件接线，再完成 DHT11、MQ135、TJC 串口屏基础调试；随后建立串口屏控件命名和事件协议，解决数据显示保持和页面闪烁问题；再加入 SmartHomeController 管理净化、新风、加湿、自动和节能状态；最后扩展 HTTP API、小程序页面和 MCP 语音控制桥接，使作品从本地显示升级为多端联动控制系统。")

    h1(doc, "第二部分  系统组成及功能说明")
    h2(doc, "2.1 整体介绍")
    add_p(doc, "系统由环境感知、控制核心、人机交互、执行反馈和网络/语音入口组成。DHT11 与 MQ135 提供环境数据；ESP32S3 对数据进行校验、评分和状态融合；串口屏、小程序和 AI/MCP 均可查看或控制系统；红蓝 LED 与舵机扇叶反馈执行结果。所有控制入口最终进入 SmartHomeController，保证净化、新风、加湿、自动、节能等状态在屏幕、小程序和语音侧一致。")
    add_doc_image(doc, interaction, 15)
    caption(doc, "图 2 三种交互方式关系图")

    h2(doc, "2.2 硬件系统介绍")
    h3(doc, "2.2.1 硬件整体介绍")
    add_p(doc, "硬件部分以 ESP32S3 开发板为核心，外接 TJC 串口屏、DHT11、MQ135、红色 LED、蓝色 LED 和 SG90 舵机扇叶。TJC 串口屏用于现场显示与触摸输入，DHT11 用于温湿度采集，MQ135 用于空气质量变化趋势采集。红色 LED 模拟净化器，亮度越高表示档位越高；蓝色 LED 模拟加湿器；SG90 舵机带动扇叶往复摆动，表示新风档位。")
    add_doc_image(doc, hw, 15)
    caption(doc, "图 3 硬件连接关系图")
    add_doc_image(doc, IMAGES["hw2"], 11.5)
    caption(doc, "图 4 整体硬件与屏幕实物照片")

    h3(doc, "2.2.2 机械设计介绍")
    add_p(doc, "新风部分采用单独扇叶插接在 SG90 舵机输出轴上的方案。当前代码按普通 0-180°舵机设计，通过不同摆动范围和步进速度表示不同档位：低档小角度慢速摆动，中档中等角度摆动，高档接近全角度快速摆动。该方案结构简单，容易观察，也便于在答辩中说明“档位变化对应实际执行动作”。")

    h3(doc, "2.2.3 电路各模块介绍")
    add_p(doc, "串口屏通过 UART2 与 ESP32S3 通信，GPIO41 接屏幕 RX，GPIO42 接屏幕 TX，波特率为 9600。DHT11 DATA 接 GPIO18，建议增加 4.7k-10k 上拉电阻。MQ135 AO 接 GPIO1/ADC1_CH0，若模块模拟输出可能超过 3.3V，需要分压保护。红色 LED 接 GPIO13，蓝色 LED 接 GPIO14，两者通过 PWM 调节亮度。舵机信号线接 GPIO21，舵机建议使用稳定 5V 供电，并与 ESP32S3 共地。")

    h2(doc, "2.3 软件系统介绍")
    h3(doc, "2.3.1 软件整体介绍")
    add_p(doc, "软件基于现有小智 ESP32 工程扩展，主要模块包括 DHT11 驱动、MQ135 驱动、SerialHmi 串口屏驱动、SmartHomeController 智能家居控制器、SmartHomeHttpServer 小程序接口和 Xiaozhi MCP Bridge。传感器任务周期采样并生成环境状态，串口屏驱动负责控件刷新和触摸事件解析，控制器统一管理设备档位和模式，HTTP API 服务小程序，MCP 桥接服务 AI 语音控制。")
    add_doc_image(doc, sw, 15)
    caption(doc, "图 5 软件数据流与控制流程图")

    h3(doc, "2.3.2 软件各模块介绍")
    add_table(doc, ["模块", "作用"], [
        ["DHT11/MQ135", "采集温湿度和空气质量原始值，异常数据不直接污染界面显示"],
        ["SerialHmi", "发送 TJC 控件命令，更新温湿度、空气等级、曲线和页面事件"],
        ["SmartHomeController", "统一管理净化、新风、加湿、自动、节能和手动环境模拟"],
        ["HTTP Server", "向小程序提供状态、历史数据、设备控制、模式控制和环境模拟接口"],
        ["MCP Bridge", "把 AI 语音工具调用转发到 ESP32 HTTP API，实现语音控制"],
    ])
    caption(doc, "表 2 软件模块说明")

    h1(doc, "第三部分  完成情况及性能参数")
    h2(doc, "3.1 整体介绍")
    add_p(doc, "当前作品已完成“采集-显示-控制-反馈-语音”的主要闭环。上电后，串口屏持续显示环境数据；用户可以切换到空气详情页查看近期空气评分曲线，也可以进入智能家居页控制设备。小程序能够在同一局域网内读取状态和历史数据，并控制三类执行器和两种模式。AI 语音控制通过 MCP 工具控制相同设备。")
    add_doc_image(doc, IMAGES["hw1"], 10.8)
    caption(doc, "图 6 整机、串口屏和小程序联调照片")

    h2(doc, "3.2 工程成果")
    h3(doc, "3.2.1 机械成果")
    add_p(doc, "完成了舵机扇叶执行机构。它可以通过档位变化改变摆动范围和摆动速度，用于直观表达新风开启和档位变化。")

    h3(doc, "3.2.2 电路成果")
    add_p(doc, "完成了 ESP32S3 与串口屏、DHT11、MQ135、红蓝 LED 和舵机的接入。各模块采用共地连接，LED 与舵机均由 LEDC PWM 输出控制。")

    h3(doc, "3.2.3 软件成果")
    add_p(doc, "串口屏页面包括首页、空气详情页、智能家居控制页和 AI/设置页。首页用于快速查看环境状态，空气详情页用于查看空气评分和趋势，智能家居页用于控制执行器，AI/设置页用于展示 AI 状态和手动环境模拟入口。")
    add_doc_image(doc, IMAGES["home"], 11.5)
    caption(doc, "图 7 串口屏首页")
    add_doc_image(doc, IMAGES["air"], 11.5)
    caption(doc, "图 8 串口屏空气详情页")
    add_doc_image(doc, IMAGES["ai"], 11.5)
    caption(doc, "图 9 串口屏 AI 与设置页")

    h2(doc, "3.3 特性成果")
    add_table(doc, ["功能", "完成情况", "应用价值"], [
        ["环境监测", "温湿度、空气原始值、空气评分和建议显示", "让用户知道当前室内状态"],
        ["串口屏控制", "页面切换、设备控制、模式控制、曲线显示", "现场不用手机也能操作"],
        ["小程序控制", "局域网读取状态、历史和控制设备", "适合坐在桌边或床边查看与控制"],
        ["AI 语音控制", "MCP 工具控制净化、新风、加湿、自动、节能", "手不方便时可直接语音操作"],
        ["自动模式", "根据温湿度和空气评分联动执行器", "从提醒升级为主动调节"],
        ["节能模式", "关闭自动和执行器", "适合离开房间或低功耗展示"],
    ])
    caption(doc, "表 3 功能完成情况与应用价值")
    add_p(doc, "需要说明的是，当前 MQ135 主要用于演示级空气质量变化和评分，不等同于专业 PM2.5、CO2 或 TVOC 浓度检测；DHT11 也属于基础温湿度传感器。作品当前重点是完成闭环架构和交互验证，后续可替换为更高精度传感器提升实际可用性。")

    h1(doc, "第四部分  总结")
    h2(doc, "4.1 可扩展之处")
    add_p(doc, "后续可扩展方向包括：接入 PM2.5、CO2、TVOC 等更专业的环境传感器；把 LED 和舵机演示执行器替换为真实净化器、风机和加湿模块；增加外壳和独立供电设计，提高作品稳定性；把历史数据保存到本地数据库或云端，形成更长时间的趋势分析；为 HTTP API 增加鉴权和配置页面，使小程序方案更接近真实产品。")

    h2(doc, "4.2 心得体会")
    add_p(doc, "本作品开发过程中，最核心的收获是把多个单独模块整合成一个稳定闭环。单独点亮 LED、读取传感器或显示串口屏并不难，真正需要反复调试的是各模块同时工作后的稳定性，例如 DHT11 时序、串口屏刷新闪烁、页面控件命名、舵机供电、HTTP API 启动时机和多入口状态一致性。通过逐步排查，我们把屏幕显示从“刷新时短暂出现”调整为持续保持，把设备控制统一到 SmartHomeController，把小程序和 AI 控制都接入同一套接口。这个过程让我们更清楚地理解了嵌入式作品不仅要“能跑”，还要“能演示、能解释、能继续扩展”。")
    add_p(doc, "从应用角度看，室内环境管理不是单纯显示几个数字，而是要把环境变化变成用户能看懂的状态、能执行的建议和能看到反馈的动作。本作品虽然仍是原型，但已经具备实际智能家居系统的基本形态：传感器负责感知，控制器负责判断，多端界面负责交互，执行器负责反馈。后续只要替换更可靠的传感器和真实执行设备，就可以继续向实用化方向推进。")

    h1(doc, "第五部分  参考文献")
    refs = [
        "[1] Espressif Systems. ESP-IDF Programming Guide.",
        "[2] Espressif Systems. ESP32-S3 Datasheet.",
        "[3] Aosong Electronics. DHT11 Humidity & Temperature Sensor Datasheet.",
        "[4] Zhengzhou Winsen Electronics. MQ135 Gas Sensor Technical Data.",
        "[5] TJC/USART HMI 串口屏开发资料.",
        "[6] Model Context Protocol Documentation.",
        "[7] 微信小程序开发文档.",
    ]
    for ref in refs:
        add_p(doc, ref, first_line=False)

    doc.save(REPORT_DOCX)


def add_slide_title(slide, title):
    box = slide.shapes.add_textbox(PptInches(0.45), PptInches(0.25), PptInches(9.1), PptInches(0.45))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = PptPt(22)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(31, 62, 93)


def ppt_textbox(slide, text, x, y, w, h, size=16, bold=False, color=(40, 40, 40), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Microsoft YaHei"
    p.font.size = PptPt(size)
    p.font.bold = bold
    p.font.color.rgb = PptRGBColor(*color)
    return box


def ppt_bullet(slide, items, x, y, w, h, size=15):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = PptPt(size)
        p.font.color.rgb = PptRGBColor(45, 45, 45)
        p.space_after = PptPt(8)
    return box


def add_picture(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), PptInches(x), PptInches(y), PptInches(w), PptInches(h))


def build_ppt(figs):
    arch, hw, sw, interaction = figs
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(5.625)
    blank = prs.slide_layouts[6]

    def bg(slide, color=(248, 250, 252)):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptRGBColor(*color)

    s = prs.slides.add_slide(blank)
    bg(s, (31, 62, 93))
    ppt_textbox(s, "喵伴空气管家", 0.6, 1.55, 8.8, 0.7, 34, True, (255, 255, 255), PP_ALIGN.CENTER)
    ppt_textbox(s, "室内环境监测与智能调节系统", 0.6, 2.35, 8.8, 0.45, 20, False, (226, 238, 248), PP_ALIGN.CENTER)
    ppt_textbox(s, "传感器 + 串口屏 + 小程序 + AI 语音 + 执行器联动", 0.6, 3.1, 8.8, 0.35, 15, False, (226, 238, 248), PP_ALIGN.CENTER)

    s = prs.slides.add_slide(blank)
    bg(s)
    add_slide_title(s, "1. 我们解决什么问题")
    ppt_bullet(s, [
        "室内空气和舒适度常常看不见，用户只能凭感觉判断",
        "普通监测器只显示数据，缺少联动控制和多端交互",
        "作品把“看见状态、理解建议、直接控制”放在同一套系统里",
    ], 0.7, 1.05, 4.6, 3.2)
    add_picture(s, IMAGES["hw2"], 5.55, 1.0, 3.75, 3.2)

    s = prs.slides.add_slide(blank)
    bg(s)
    add_slide_title(s, "2. 系统总体架构")
    add_picture(s, arch, 0.55, 0.85, 8.9, 4.35)

    s = prs.slides.add_slide(blank)
    bg(s)
    add_slide_title(s, "3. 硬件与执行反馈")
    add_picture(s, hw, 0.45, 0.85, 5.35, 3.0)
    ppt_bullet(s, [
        "DHT11：温湿度采集",
        "MQ135：空气质量原始值与评分",
        "红 LED：净化档位，蓝 LED：加湿档位",
        "SG90 舵机扇叶：新风档位，0-180° 往复摆动",
    ], 6.0, 1.05, 3.55, 3.6, 14)

    s = prs.slides.add_slide(blank)
    bg(s)
    add_slide_title(s, "4. 串口屏：现场可见、现场可控")
    add_picture(s, IMAGES["home"], 0.55, 1.0, 2.85, 1.62)
    add_picture(s, IMAGES["air"], 3.62, 1.0, 2.85, 1.62)
    add_picture(s, IMAGES["ai"], 6.69, 1.0, 2.85, 1.62)
    ppt_bullet(s, [
        "首页显示温湿度和空气状态",
        "空气详情页显示空气评分、MQ135 原始值、舒适度和曲线",
        "智能家居页控制净化、新风、加湿、自动和节能",
        "页面保持上一轮数据，新数据到来时覆盖更新，减少闪烁",
    ], 0.75, 3.0, 8.8, 1.7, 14)

    s = prs.slides.add_slide(blank)
    bg(s)
    add_slide_title(s, "5. 小程序与 AI 语音控制")
    add_picture(s, interaction, 0.55, 0.85, 5.6, 3.0)
    ppt_bullet(s, [
        "小程序通过局域网 HTTP API 读取状态、历史数据和控制设备",
        "AI 语音通过 MCP 工具调用桥接到 ESP32 HTTP API",
        "三个入口都进入同一控制器，保证显示、状态和执行结果一致",
    ], 6.35, 1.05, 3.1, 3.55, 14)

    s = prs.slides.add_slide(blank)
    bg(s)
    add_slide_title(s, "6. 自动联动与应用价值")
    ppt_bullet(s, [
        "污染场景：自动提高净化和新风档位",
        "干燥场景：自动开启加湿",
        "节能场景：关闭自动控制和所有执行器",
        "适合宿舍、家庭、办公室、教室等小型空间",
        "后续替换专业传感器和真实执行设备即可向实用化扩展",
    ], 0.85, 1.0, 4.7, 3.7, 15)
    add_picture(s, IMAGES["hw1"], 5.85, 1.0, 3.4, 3.85)

    s = prs.slides.add_slide(blank)
    bg(s, (31, 62, 93))
    ppt_textbox(s, "总结", 0.6, 0.8, 8.8, 0.55, 28, True, (255, 255, 255), PP_ALIGN.CENTER)
    ppt_textbox(s, "我们完成的不是单个传感器演示，而是一个能感知、能显示、能控制、能语音交互的室内环境管理原型。", 0.9, 1.8, 8.2, 1.2, 22, False, (255, 255, 255), PP_ALIGN.CENTER)
    ppt_textbox(s, "下一步：更高精度传感器、真实执行设备、外壳与长期数据记录。", 0.9, 3.45, 8.2, 0.55, 16, False, (226, 238, 248), PP_ALIGN.CENTER)

    prs.save(PPTX_OUT)


def build_pdf(figs):
    arch, hw, sw, interaction = figs
    font_path = "C:/Windows/Fonts/simsun.ttc"
    bold_path = "C:/Windows/Fonts/simhei.ttf"
    pdfmetrics.registerFont(TTFont("SimSun", font_path))
    pdfmetrics.registerFont(TTFont("SimHei", bold_path))
    styles = getSampleStyleSheet()
    normal = ParagraphStyle("cn", parent=styles["Normal"], fontName="SimSun", fontSize=10.5, leading=17, firstLineIndent=21)
    h = ParagraphStyle("h", parent=styles["Heading1"], fontName="SimHei", fontSize=16, leading=24, alignment=TA_CENTER, spaceBefore=10, spaceAfter=8)
    h2s = ParagraphStyle("h2", parent=styles["Heading2"], fontName="SimHei", fontSize=13, leading=20, spaceBefore=8, spaceAfter=5)
    cap = ParagraphStyle("cap", parent=styles["Normal"], fontName="SimSun", fontSize=9, leading=13, alignment=TA_CENTER)
    def img(path, w=15, h=8.2):
        return PdfImage(str(path), width=w * cm, height=h * cm)

    metric_table = Table(
        [
            ["指标项", "当前实现"],
            ["主控芯片", "ESP32S3，基于现有小智 ESP32 工程扩展"],
            ["传感器", "DHT11 温湿度；MQ135 模拟空气质量原始值"],
            ["数据刷新", "传感器任务约 5 秒采样一次，屏幕保持上一帧有效数据"],
            ["串口屏通信", "UART2，GPIO41/GPIO42，9600bps，8N1"],
            ["执行器", "GPIO13 红 LED、GPIO14 蓝 LED、GPIO21 SG90 舵机扇叶"],
            ["小程序接口", "状态、历史、设备、模式和环境模拟接口"],
        ],
        colWidths=[4 * cm, 11 * cm],
        style=TableStyle([("FONT", (0, 0), (-1, -1), "SimSun"), ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9EAF7")), ("GRID", (0, 0), (-1, -1), 0.5, colors.grey), ("VALIGN", (0, 0), (-1, -1), "MIDDLE")]),
    )

    function_table = Table(
        [
            ["功能", "完成情况", "应用价值"],
            ["环境监测", "温湿度、空气评分、建议显示", "让用户知道当前室内状态"],
            ["串口屏控制", "页面切换、设备控制、曲线显示", "现场不用手机也能操作"],
            ["小程序控制", "局域网读取状态与控制设备", "便于远程查看和调档"],
            ["AI 语音控制", "MCP 工具控制设备和模式", "适合免手操作"],
            ["自动/节能", "按环境联动或关闭执行器", "从提醒升级为主动调节"],
        ],
        colWidths=[3 * cm, 6 * cm, 6 * cm],
        style=TableStyle([("FONT", (0, 0), (-1, -1), "SimSun"), ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9EAF7")), ("GRID", (0, 0), (-1, -1), 0.5, colors.grey), ("VALIGN", (0, 0), (-1, -1), "MIDDLE")]),
    )

    story = [
        Paragraph("喵伴空气管家作品设计报告", h),
        Paragraph("室内环境监测与智能调节系统", cap),
        Spacer(1, 0.4 * cm),
        Paragraph("摘 要", h2s),
        Paragraph("本作品面向宿舍、家庭、办公室等小型室内空间，设计并实现了一套“能监测、能显示、能控制、能语音交互”的室内环境智能调节系统。系统以 ESP32S3 为控制核心，接入 DHT11 温湿度传感器和 MQ135 空气质量传感器，使用 TJC 串口屏持续显示温度、湿度、空气等级、空气评分和近期变化曲线；同时通过红色 LED、蓝色 LED 和 SG90 舵机扇叶分别模拟净化、加湿和新风执行机构。用户可以在串口屏上直接控制设备，也可以通过微信小程序在同一局域网内查看状态、调节档位和查看历史数据，还可以通过 AI 语音/MCP 工具完成免手操作。", normal),
        Paragraph("关键词：ESP32S3；室内空气监测；串口屏；微信小程序；MCP 语音控制；智能家居", normal),
        Paragraph("第一部分 作品概述", h),
        Paragraph("1.1 功能与特性", h2s),
        Paragraph("作品完成了室内温湿度与空气质量采集、串口屏可视化显示、屏幕触摸控制、小程序局域网控制、AI 语音控制和自动联动调节等功能。串口屏首页展示温度、湿度和空气状态，空气详情页展示空气评分、MQ135 原始值、舒适度建议和近期空气质量曲线，智能家居页提供净化、新风、加湿、自动、节能等按键。", normal),
        img(arch),
        Paragraph("图 1 系统总体架构图", cap),
        Paragraph("1.2 应用领域", h2s),
        Paragraph("本作品适用于宿舍、家庭、办公室、教室、实验室等小型室内空间的环境监测与智能联动演示。对于实际应用，它可以提醒用户关注温湿度和空气状态，并在空气较差或环境干燥时自动给出调节建议；对于课程设计和竞赛展示，它把传感器、串口屏、网络接口、语音工具和执行器整合在一起。", normal),
        Paragraph("1.3 主要技术特点", h2s),
        Paragraph("系统以 ESP32S3 为统一控制核心，完成传感器采集、串口屏通信、HTTP API、MCP 工具注册和执行器控制；TJC 串口屏负责静态页面和控件显示，固件只更新动态数据，减少整屏重绘；小程序采用局域网 HTTP 方案；AI 语音控制通过 MCP 工具把语音意图转为设备控制命令。", normal),
        Paragraph("1.4 主要性能指标", h2s),
        metric_table,
        Paragraph("第二部分 系统组成及功能说明", h),
        Paragraph("2.1 整体介绍", h2s),
        Paragraph("系统由环境感知、控制核心、人机交互、执行反馈和网络/语音入口组成。DHT11 与 MQ135 提供环境数据；ESP32S3 对数据进行校验、评分和状态融合；串口屏、小程序和 AI/MCP 均可查看或控制系统；红蓝 LED 与舵机扇叶反馈执行结果。", normal),
        img(interaction),
        Paragraph("图 2 三种交互方式关系图", cap),
        Paragraph("2.2 硬件系统介绍", h2s),
        Paragraph("硬件部分以 ESP32S3 开发板为核心，外接 TJC 串口屏、DHT11、MQ135、红色 LED、蓝色 LED 和 SG90 舵机扇叶。红色 LED 模拟净化器，亮度越高表示档位越高；蓝色 LED 模拟加湿器；SG90 舵机带动扇叶往复摆动，表示新风档位。", normal),
        img(hw),
        Paragraph("图 3 硬件连接关系图", cap),
        img(IMAGES["hw2"], 11.5, 8.6),
        Paragraph("图 4 整体硬件与屏幕实物照片", cap),
        Paragraph("新风部分采用单独扇叶插接在 SG90 舵机输出轴上的方案。当前代码按普通 0-180°舵机设计，通过不同摆动范围和步进速度表示不同档位。串口屏通过 UART2 与 ESP32S3 通信，GPIO41 接屏幕 RX，GPIO42 接屏幕 TX；DHT11 接 GPIO18；MQ135 接 GPIO1/ADC1_CH0；红蓝 LED 分别接 GPIO13 和 GPIO14；舵机信号线接 GPIO21。", normal),
        Paragraph("2.3 软件系统介绍", h2s),
        Paragraph("软件基于现有小智 ESP32 工程扩展，主要模块包括 DHT11 驱动、MQ135 驱动、SerialHmi 串口屏驱动、SmartHomeController 智能家居控制器、SmartHomeHttpServer 小程序接口和 Xiaozhi MCP Bridge。", normal),
        img(sw),
        Paragraph("图 5 软件数据流与控制流程图", cap),
        Paragraph("第三部分 完成情况及性能参数", h),
        Paragraph("3.1 整体介绍", h2s),
        Paragraph("当前作品已完成“采集-显示-控制-反馈-语音”的主要闭环。上电后，串口屏持续显示环境数据；用户可以切换到空气详情页查看近期空气评分曲线，也可以进入智能家居页控制设备。小程序能够在同一局域网内读取状态和历史数据，并控制三类执行器和两种模式。AI 语音控制通过 MCP 工具控制相同设备。", normal),
        img(IMAGES["hw1"], 10.5, 13.9),
        Paragraph("图 6 整机、串口屏和小程序联调照片", cap),
        PageBreak(),
        Paragraph("3.2 工程成果", h2s),
        Paragraph("串口屏页面包括首页、空气详情页、智能家居控制页和 AI/设置页。首页用于快速查看环境状态，空气详情页用于查看空气评分和趋势，智能家居页用于控制执行器，AI/设置页用于展示 AI 状态和手动环境模拟入口。", normal),
        img(IMAGES["home"], 12.5, 7.1),
        Paragraph("图 7 串口屏首页", cap),
        img(IMAGES["air"], 12.5, 7.1),
        Paragraph("图 8 串口屏空气详情页", cap),
        img(IMAGES["ai"], 12.5, 7.1),
        Paragraph("图 9 串口屏 AI 与设置页", cap),
        Paragraph("3.3 特性成果", h2s),
        function_table,
        Spacer(1, 0.3 * cm),
        Paragraph("需要说明的是，当前 MQ135 主要用于演示级空气质量变化和评分，不等同于专业 PM2.5、CO2 或 TVOC 浓度检测；DHT11 也属于基础温湿度传感器。作品当前重点是完成闭环架构和交互验证，后续可替换为更高精度传感器提升实际可用性。", normal),
        Paragraph("第四部分 总结", h),
        Paragraph("4.1 可扩展之处", h2s),
        Paragraph("后续可扩展方向包括：接入 PM2.5、CO2、TVOC 等更专业的环境传感器；把 LED 和舵机演示执行器替换为真实净化器、风机和加湿模块；增加外壳和独立供电设计，提高作品稳定性；把历史数据保存到本地数据库或云端；为 HTTP API 增加鉴权和配置页面。", normal),
        Paragraph("4.2 心得体会", h2s),
        Paragraph("本作品开发过程中，最核心的收获是把多个单独模块整合成一个稳定闭环。单独点亮 LED、读取传感器或显示串口屏并不难，真正需要反复调试的是各模块同时工作后的稳定性，例如 DHT11 时序、串口屏刷新闪烁、页面控件命名、舵机供电、HTTP API 启动时机和多入口状态一致性。通过逐步排查，我们把屏幕显示从“刷新时短暂出现”调整为持续保持，把设备控制统一到 SmartHomeController，把小程序和 AI 控制都接入同一套接口。", normal),
        Paragraph("第五部分 参考文献", h),
        Paragraph("[1] Espressif Systems. ESP-IDF Programming Guide. [2] Espressif Systems. ESP32-S3 Datasheet. [3] Aosong Electronics. DHT11 Humidity & Temperature Sensor Datasheet. [4] Zhengzhou Winsen Electronics. MQ135 Gas Sensor Technical Data. [5] TJC/USART HMI 串口屏开发资料. [6] Model Context Protocol Documentation. [7] 微信小程序开发文档.", normal),
    ]
    doc = SimpleDocTemplate(str(REPORT_PDF), pagesize=A4, rightMargin=2 * cm, leftMargin=2.5 * cm, topMargin=2.5 * cm, bottomMargin=2 * cm)
    doc.build(story)


def build_guide():
    md = """# 视频讲稿与答辩建议

## 3 分钟作品视频讲稿

0:00-0:15 开场：大家好，我们的作品叫“喵伴空气管家”，面向宿舍、家庭和办公室等小型室内空间，解决空气状态看不见、调节不直观、控制入口分散的问题。

0:15-0:40 硬件介绍：系统以 ESP32S3 为核心，接入 DHT11 采集温湿度，MQ135 采集空气质量原始值，TJC 串口屏负责现场显示和触摸控制。右侧红色 LED 表示净化器，蓝色 LED 表示加湿器，舵机扇叶表示新风档位。

0:40-1:20 串口屏演示：首页持续显示温度、湿度和空气状态；空气详情页展示空气评分、原始值、舒适度建议和近期空气质量曲线；智能家居页可以控制净化、新风、加湿、自动和节能。数据不是闪一下就消失，而是保持上一轮结果，新数据到来时覆盖更新。

1:20-1:55 小程序演示：小程序通过同一局域网访问 ESP32S3，可以查看当前状态、历史数据，也可以调节净化、新风和加湿档位。这个入口适合用户不在设备旁边，但仍在同一房间或同一网络下查看和控制。

1:55-2:30 AI 语音演示：AI 语音通过 MCP 工具调用，把“打开净化器”“新风调到二档”“开启节能模式”等指令转成接口请求，发送给 ESP32S3。语音、小程序和串口屏最终进入同一套控制器，所以状态不会互相打架。

2:30-2:50 自动联动演示：在污染或干燥场景下，自动模式会根据空气评分、温湿度调整净化、新风和加湿档位；节能模式会关闭自动和执行器，适合离开房间时使用。

2:50-3:00 收尾：这套系统目前是低成本原型，已经完成感知、显示、控制、语音和自动联动闭环。后续替换高精度传感器和真实执行设备，就可以向实际室内环境管理产品扩展。

## 5 分钟答辩安排

视频建议控制在 2 分 40 秒到 2 分 55 秒，PPT 只讲 1 分 30 秒到 2 分钟。PPT 不要逐字念，按“问题-方案-完成情况-价值-不足与改进”讲。

1. 第 1 页，10 秒：报作品名和一句话定位。
2. 第 2 页，20 秒：说清楚应用问题，不要陷入器件细节。
3. 第 3 页，25 秒：讲总体架构，强调三种入口最终统一控制。
4. 第 4 页，20 秒：讲硬件和执行反馈，让评委知道每个模块做什么。
5. 第 5 页，25 秒：讲串口屏完成了持续显示、曲线和控制。
6. 第 6 页，25 秒：讲小程序和 AI 语音的应用价值。
7. 第 7-8 页，25 秒：总结自动联动、应用场景和后续改进。

## 视频拍摄建议

先拍“问题场景”：桌面上放作品，旁白说室内环境需要看得见也需要能控制。

再拍“硬件全景”：镜头从 ESP32S3、传感器、串口屏、LED、舵机扇叶扫过，停留在整体结构上。

串口屏演示要稳：用三脚架或固定手机，先拍首页，再点到空气详情页，再点到智能家居页。每个页面停 3-5 秒，让评委看清字。

执行器演示要近景：净化调到 1/2/3 档时拍红 LED 亮度变化，加湿拍蓝 LED，新风拍舵机扇叶摆动。最好用同一个镜头连续拍，证明是实时控制。

小程序演示要把手机和实物放同一画面：点击小程序按钮后，镜头里同时看到 LED 或舵机变化，这比单独录屏更有说服力。

AI 演示要短：只保留 2-3 条最稳的语音指令，例如“打开净化器二档”“开启节能模式”“查询当前空气状态”。不要现场临时发挥长句。

最后用一个应用场景收束：宿舍睡前看空气状态、干燥时自动加湿、空气差时净化和新风联动。评委要听到“为什么有用”，不只是“我能控制 GPIO”。

## 可能被问到的问题

Q1：MQ135 能不能准确测 PM2.5 或 CO2？
A：当前 MQ135 用于演示级空气质量趋势和评分，不直接等同于 PM2.5、CO2 或 TVOC 的精确浓度。真实产品会换成专业传感器并做标定。

Q2：为什么要同时做串口屏、小程序和 AI？
A：三种入口对应不同使用场景：串口屏适合设备旁现场操作，小程序适合同一局域网内远程查看，AI 适合免手控制。它们最终进入同一控制器，所以状态统一。

Q3：自动模式的规则是什么？
A：空气评分低或 MQ135 原始值异常时提高净化和新风档位，湿度低时开启加湿，高温或需要通风时提高新风；节能模式会关闭自动和所有执行器。

Q4：小程序为什么使用局域网 HTTP，不上云？
A：答辩阶段优先保证现场稳定和可解释性。同一 WiFi 下直接访问 ESP32S3，链路短、调试快。后续可加云端、鉴权和长期数据存储。

Q5：如果网络断了还能用吗？
A：串口屏本地显示和本地按键控制仍可用；小程序和 AI 语音依赖网络或桥接服务。这个设计保证基础功能不完全依赖手机或云端。

Q6：作品和普通空气检测仪相比有什么价值？
A：普通检测仪主要显示数据，本作品增加了执行器反馈、自动联动、小程序和 AI 语音控制，展示的是从感知到控制的闭环。

Q7：当前最大的不足是什么？
A：传感器精度和执行器还处于原型演示级，外壳、电源、安全和长期数据记录还需要完善。我们已经把接口和控制架构搭好，后续替换硬件即可扩展。
"""
    GUIDE_MD.write_text(md, encoding="utf-8")

    doc = Document()
    style_document(doc)
    for line in md.splitlines():
        if line.startswith("# "):
            h1(doc, line[2:])
        elif line.startswith("## "):
            h2(doc, line[3:])
        elif line.strip():
            add_p(doc, line.strip(), first_line=not line.startswith(("0:", "1.", "2.", "3.", "4.", "5.", "6.", "7.", "Q", "A")))
    doc.save(GUIDE_DOCX)


def main():
    figs = make_figures()
    build_report(figs)
    build_pdf(figs)
    build_ppt(figs)
    build_guide()
    for path in [REPORT_DOCX, REPORT_PDF, PPTX_OUT, GUIDE_MD, GUIDE_DOCX]:
        print(path)


if __name__ == "__main__":
    main()
