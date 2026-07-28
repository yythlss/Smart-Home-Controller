# -*- coding: utf-8 -*-
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "report_output_rebuild"
ASSET_DIR = OUT / "assets"
FINAL_DIR = OUT / "final"
TMP_DIR = OUT / "tmp" / "ppt_opt_assets"

PPTX_OPT = FINAL_DIR / "喵伴空气管家答辩PPT_优化版_2026-07-08.pptx"

GREEN = "174238"
ORANGE = "D9822B"
LIGHT = "F5F6F3"

FONT = "Microsoft YaHei"
FONT_BOLD = "Microsoft YaHei UI"


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def ensure_dirs() -> None:
    FINAL_DIR.mkdir(parents=True, exist_ok=True)
    TMP_DIR.mkdir(parents=True, exist_ok=True)


def crop_box(src: Path, dst_name: str, box: tuple[float, float, float, float]) -> Path:
    dst = TMP_DIR / dst_name
    with Image.open(src) as im:
        w, h = im.size
        left = int(w * box[0])
        top = int(h * box[1])
        right = int(w * box[2])
        bottom = int(h * box[3])
        cropped = im.crop((left, top, right, bottom))
        cropped.save(dst, quality=92)
    return dst


def crop_to_aspect(src: Path, dst_name: str, aspect: float) -> Path:
    dst = TMP_DIR / dst_name
    with Image.open(src) as im:
        w, h = im.size
        current = w / h
        if current > aspect:
            new_w = int(h * aspect)
            left = (w - new_w) // 2
            box = (left, 0, left + new_w, h)
        else:
            new_h = int(w / aspect)
            top = (h - new_h) // 2
            box = (0, top, w, top + new_h)
        im.crop(box).save(dst, quality=92)
    return dst


def set_bg(slide, color: str = LIGHT) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 18,
    bold: bool = False,
    color: str = GREEN,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT_BOLD if bold else FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)


def add_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    color = LIGHT if dark else GREEN
    add_text(slide, title, 0.55, 0.38, 8.9, 0.55, 28, True, color)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.98, 10.8, 0.34, 13, False, color)


def add_footer(slide, index: int) -> None:
    add_text(slide, f"{index:02d}", 12.32, 6.96, 0.45, 0.2, 10, True, ORANGE, PP_ALIGN.RIGHT)


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: str = ORANGE, width: float = 2.0) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)


def add_round_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = LIGHT,
    line: str = GREEN,
    width: float = 1.2,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)


def add_chip(slide, text: str, x: float, y: float, w: float, color: str = GREEN, fill: str = LIGHT) -> None:
    add_round_rect(slide, x, y, w, 0.44, fill=fill, line=color, width=1.0)
    add_text(slide, text, x + 0.12, y + 0.09, w - 0.24, 0.18, 10, True, color, PP_ALIGN.CENTER)


def add_picture(slide, path: Path, x: float, y: float, w: float | None = None, h: float | None = None) -> None:
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def add_picture_frame(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    crop = crop_to_aspect(path, f"crop_{path.stem}_{int(w*100)}_{int(h*100)}.jpg", w / h)
    add_picture(slide, crop, x, y, w=w, h=h)
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    rect.fill.background()
    rect.line.color.rgb = rgb(GREEN)
    rect.line.width = Pt(1.2)


def slide_title(prs: Presentation, hero: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, GREEN)
    add_picture_frame(slide, hero, 6.8, 0.55, 5.85, 6.38)
    add_text(slide, "喵伴空气管家", 0.72, 1.05, 5.5, 0.7, 38, True, LIGHT)
    add_text(slide, "室内环境智能调节系统", 0.75, 1.85, 4.8, 0.36, 18, False, LIGHT)
    add_text(slide, "把温湿度、空气质量、屏幕控制、手机控制和语音控制接到同一套状态里。", 0.76, 3.15, 5.25, 0.85, 21, True, LIGHT)
    add_chip(slide, "串口屏", 0.78, 5.55, 1.0, ORANGE, GREEN)
    add_chip(slide, "小程序", 1.95, 5.55, 1.0, ORANGE, GREEN)
    add_chip(slide, "AI 语音", 3.12, 5.55, 1.12, ORANGE, GREEN)
    add_chip(slide, "自动联动", 4.45, 5.55, 1.18, ORANGE, GREEN)
    add_footer(slide, 1)


def slide_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "先解决一个很具体的问题", "室内环境不是只要“能测到”，还要让人看得懂、控得住，异常时能自动处理。")
    items = [
        ("01", "看不直观", "闷、干、热、空气差，很多时候是凭感觉发现。"),
        ("02", "入口分散", "屏幕、手机、语音如果不同步，控制结果容易混乱。"),
        ("03", "手动容易忘", "需要自动模式处理明显异常，也需要节能模式一键收口。"),
    ]
    for idx, (num, head, body) in enumerate(items):
        y = 1.65 + idx * 1.45
        add_round_rect(slide, 0.82, y, 11.5, 1.05, fill=LIGHT, line=GREEN)
        add_text(slide, num, 1.08, y + 0.28, 0.65, 0.3, 18, True, ORANGE, PP_ALIGN.CENTER)
        add_line(slide, 1.9, y + 0.17, 1.9, y + 0.88, ORANGE, 2.4)
        add_text(slide, head, 2.18, y + 0.18, 2.1, 0.25, 17, True, GREEN)
        add_text(slide, body, 4.3, y + 0.22, 6.7, 0.3, 17, False, GREEN)
    add_text(slide, "我们的思路：不是做一个单独的读数页面，而是把“检测 - 判断 - 显示 - 控制”连成一条线。", 1.0, 6.25, 11.0, 0.35, 18, True, ORANGE, PP_ALIGN.CENTER)
    add_footer(slide, 2)


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "系统主线：所有入口共用一份状态", "传感器进 ESP32S3，屏幕、小程序和语音都读写同一套设备状态。")
    blocks = [
        (0.72, "感知", "DHT11\nMQ135"),
        (3.28, "ESP32S3", "采样\n评分\n状态缓存\nHTTP API"),
        (6.13, "交互入口", "串口屏\n小程序\nAI/MCP"),
        (9.0, "执行", "净化红灯\n新风舵机\n加湿蓝灯"),
    ]
    for x, head, body in blocks:
        add_round_rect(slide, x, 1.72, 2.2, 3.58, fill=LIGHT, line=GREEN, width=1.5)
        add_text(slide, head, x + 0.25, 2.05, 1.7, 0.28, 18, True, GREEN, PP_ALIGN.CENTER)
        add_line(slide, x + 0.35, 2.52, x + 1.85, 2.52, ORANGE, 2.2)
        add_text(slide, body, x + 0.25, 2.88, 1.7, 1.2, 17, False, GREEN, PP_ALIGN.CENTER)
    for x in [2.95, 5.77, 8.62]:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(3.15), Inches(0.72), Inches(0.42))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = rgb(ORANGE)
        arrow.line.fill.background()
    add_text(slide, "这条主线保证：从屏幕点按钮、从手机点档位、从语音下指令，最后改的是同一组状态。", 0.95, 6.15, 11.5, 0.35, 18, True, GREEN, PP_ALIGN.CENTER)
    add_footer(slide, 3)


def slide_screen(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "串口屏：现场能看，也能直接控制", "屏幕不是摆设，现场不拿手机也能完成查看和操作。")
    shots = [
        ("首页", ASSET_DIR / "serial_page_home.png", "温度、湿度、空气评分"),
        ("空气详情", ASSET_DIR / "serial_page_air_score.png", "MQ135 原始值和曲线"),
        ("AI/设置", ASSET_DIR / "serial_page_ai_settings.png", "手动环境场景"),
    ]
    for idx, (name, path, note) in enumerate(shots):
        x = 0.72 + idx * 4.18
        add_picture_frame(slide, path, x, 1.55, 3.5, 2.05)
        add_text(slide, name, x, 3.9, 3.5, 0.28, 17, True, GREEN, PP_ALIGN.CENTER)
        add_text(slide, note, x, 4.35, 3.5, 0.28, 13, False, GREEN, PP_ALIGN.CENTER)
    add_round_rect(slide, 1.12, 5.5, 10.9, 0.72, fill=LIGHT, line=ORANGE, width=1.4)
    add_text(slide, "显示策略：数据保持在屏幕上，新数据来了再覆盖旧数据，避免“刷新一下就消失”。", 1.35, 5.76, 10.4, 0.22, 17, True, GREEN, PP_ALIGN.CENTER)
    add_footer(slide, 4)


def slide_mini_program(prs: Presentation, phone_crop: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "小程序：手机端看状态、调档位", "同一 WiFi 下，手机直接访问 ESP32S3 的 HTTP API。")
    add_picture_frame(slide, phone_crop, 0.8, 1.25, 4.0, 5.25)
    api_items = [
        ("GET /api/state", "读取当前温湿度、空气评分、设备档位"),
        ("POST /api/device", "净化、新风、加湿切换 0-3 档"),
        ("POST /api/mode", "自动模式、节能模式开关"),
        ("GET /api/history", "查看最近 30 条空气评分"),
    ]
    for i, (api, desc) in enumerate(api_items):
        y = 1.55 + i * 1.05
        add_text(slide, api, 5.55, y, 2.55, 0.25, 15, True, ORANGE)
        add_text(slide, desc, 8.25, y, 3.95, 0.25, 15, False, GREEN)
        add_line(slide, 5.55, y + 0.42, 12.15, y + 0.42, GREEN, 0.8)
    add_text(slide, "好处很直接：不用另搭服务器，答辩现场和实际同网段使用都能跑通。", 5.55, 6.15, 6.75, 0.3, 17, True, GREEN)
    add_footer(slide, 5)


def slide_ai(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "AI 语音：把一句话变成一次工具调用", "AI 负责理解口语，真正的设备状态仍由 ESP32S3 控制器统一维护。")
    labels = [
        ("用户语音", "打开净化器二档"),
        ("小智 MCP", "选择 home_set_purifier"),
        ("桥接脚本", "转发到 /api/device"),
        ("ESP32S3", "GPIO13 PWM 变化"),
    ]
    for i, (head, body) in enumerate(labels):
        x = 0.65 + i * 3.1
        add_round_rect(slide, x, 2.0, 2.42, 1.55, fill=LIGHT, line=GREEN)
        add_text(slide, head, x + 0.18, 2.25, 2.05, 0.22, 15, True, GREEN, PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.18, 2.78, 2.05, 0.28, 13, False, GREEN, PP_ALIGN.CENTER)
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x + 2.45), Inches(2.55), Inches(0.48), Inches(0.34))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(ORANGE)
            arrow.line.fill.background()
    examples = [
        "打开新风三档",
        "进入节能模式",
        "模拟干燥环境",
    ]
    add_text(slide, "答辩时可以这样说", 1.0, 4.55, 2.4, 0.25, 17, True, GREEN)
    for i, item in enumerate(examples):
        add_chip(slide, item, 3.3 + i * 2.25, 4.48, 1.75, ORANGE, LIGHT)
    add_text(slide, "重点不是“AI 直接接管硬件”，而是语音入口和手机、屏幕共用一套控制链路。", 1.0, 5.82, 11.2, 0.35, 18, True, GREEN, PP_ALIGN.CENTER)
    add_footer(slide, 6)


def slide_auto_eco(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "自动与节能：规则清楚，动作看得见", "自动模式根据环境状态调整档位，节能模式负责快速关闭执行器。")
    rows = [
        ("空气评分低", "净化升档 + 新风摆动"),
        ("湿度偏低", "加湿升档"),
        ("高温场景", "提醒通风，配合新风动作"),
        ("节能开启", "关闭自动 + 三个执行器归零"),
    ]
    add_text(slide, "环境状态", 1.0, 1.55, 2.2, 0.24, 15, True, ORANGE)
    add_text(slide, "系统动作", 6.75, 1.55, 2.2, 0.24, 15, True, ORANGE)
    for i, (left, right) in enumerate(rows):
        y = 2.0 + i * 0.9
        add_round_rect(slide, 0.95, y, 3.7, 0.56, fill=LIGHT, line=GREEN)
        add_round_rect(slide, 7.0, y, 4.0, 0.56, fill=LIGHT, line=GREEN)
        add_text(slide, left, 1.15, y + 0.16, 3.25, 0.16, 13, True, GREEN, PP_ALIGN.CENTER)
        add_text(slide, right, 7.2, y + 0.16, 3.6, 0.16, 13, True, GREEN, PP_ALIGN.CENTER)
        add_line(slide, 4.75, y + 0.28, 6.83, y + 0.28, ORANGE, 2.0)
    add_text(slide, "现场演示用“舒适 / 高温 / 干燥 / 污染”四个场景切换，不用真的等环境变差。", 1.0, 6.05, 11.2, 0.32, 18, True, GREEN, PP_ALIGN.CENTER)
    add_footer(slide, 7)


def slide_verified(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "完成情况：每条链路都有可观察结果", "每项功能都能通过屏幕、手机、语音入口、日志或外设动作确认。")
    items = [
        ("传感器", "温湿度和空气评分进入状态缓存"),
        ("串口屏", "首页、详情、AI 设置页持续显示"),
        ("执行器", "红灯、蓝灯、舵机按档位动作"),
        ("小程序", "状态读取、历史读取、设备控制"),
        ("AI/MCP", "语音工具调用转发到 HTTP API"),
        ("模式", "自动联动、节能关闭执行器"),
    ]
    for i, (head, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.85 + col * 6.0
        y = 1.65 + row * 1.35
        add_round_rect(slide, x, y, 5.25, 0.88, fill=LIGHT, line=GREEN)
        add_text(slide, head, x + 0.22, y + 0.2, 1.3, 0.2, 14, True, ORANGE)
        add_text(slide, body, x + 1.55, y + 0.2, 3.4, 0.2, 14, False, GREEN)
    add_text(slide, "演示路径：屏幕查看状态 → 小程序切换档位 → 语音控制 → 自动模式联动。", 1.0, 6.22, 11.2, 0.3, 17, True, GREEN, PP_ALIGN.CENTER)
    add_footer(slide, 8)


def slide_value(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, GREEN)
    add_title(slide, "应用价值：小空间也能做成可视化空气管理", dark=True)
    scenarios = ["家庭", "宿舍", "教室", "小办公室"]
    for i, name in enumerate(scenarios):
        x = 0.9 + i * 3.05
        add_round_rect(slide, x, 1.65, 2.35, 1.1, fill=GREEN, line=LIGHT, width=1.4)
        add_text(slide, name, x + 0.15, 1.98, 2.05, 0.24, 18, True, LIGHT, PP_ALIGN.CENTER)
    add_text(slide, "当前版本把“看见状态、手动控制、语音控制、自动联动”放在同一台 ESP32S3 上，适合低成本部署和后续扩展。", 1.0, 3.55, 11.3, 0.7, 25, True, LIGHT, PP_ALIGN.CENTER)
    add_text(slide, "下一步可以接入云端记录、PM2.5/CO2/TVOC 传感器，以及真实净化器、风机和加湿器。", 1.2, 5.35, 10.8, 0.36, 18, False, LIGHT, PP_ALIGN.CENTER)
    add_footer(slide, 9)


def generate() -> Path:
    ensure_dirs()
    hero = crop_to_aspect(ASSET_DIR / "hardware_overview_1.jpg", "hero_16_9.jpg", 0.92)
    phone_crop = crop_box(ASSET_DIR / "hardware_overview_1.jpg", "phone_crop.jpg", (0.00, 0.43, 0.53, 1.00))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_title(prs, hero)
    slide_problem(prs)
    slide_architecture(prs)
    slide_screen(prs)
    slide_mini_program(prs, phone_crop)
    slide_ai(prs)
    slide_auto_eco(prs)
    slide_verified(prs)
    slide_value(prs)
    prs.save(PPTX_OPT)
    return PPTX_OPT


if __name__ == "__main__":
    print(generate())
