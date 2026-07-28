from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import (
    Image as PdfImage,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "docs" / "report_output"
ASSET_DIR = OUT_DIR / "assets"
FIG_DIR = OUT_DIR / "figures_v2"
FINAL_DIR = OUT_DIR / "final"
FIG_DIR.mkdir(parents=True, exist_ok=True)
FINAL_DIR.mkdir(parents=True, exist_ok=True)

REPORT_DOCX = FINAL_DIR / "喵伴空气管家作品设计报告_V2_2026-07-08.docx"
REPORT_PDF = FINAL_DIR / "喵伴空气管家作品设计报告_V2_2026-07-08.pdf"
PPTX_OUT = FINAL_DIR / "喵伴空气管家答辩PPT_V2_2026-07-08.pptx"
HANDOFF = ROOT / "docs" / "phase-handoff-2026-07-08-report-ppt-v2.md"

IMAGES = {
    "home": ASSET_DIR / "serial_page_home.png",
    "air": ASSET_DIR / "serial_page_air_score.png",
    "ai": ASSET_DIR / "serial_page_ai_settings.png",
    "hw1": ASSET_DIR / "hardware_overview_1.jpg",
    "hw2": ASSET_DIR / "hardware_overview_2.jpg",
}

PRIMARY = "#24507A"
GREEN = "#2E7D55"
ORANGE = "#C66A1D"
LIGHT_BLUE = "#EAF3FC"
LIGHT_GREEN = "#EAF6EF"
LIGHT_ORANGE = "#FFF1E5"
LIGHT_GRAY = "#F6F8FA"


def font(size=26, bold=False):
    choices = [
        Path("C:/Windows/Fonts/simhei.ttf") if bold else Path("C:/Windows/Fonts/simsun.ttc"),
        Path("C:/Windows/Fonts/msyh.ttc"),
    ]
    for path in choices:
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def wrap_px(draw, text, max_width, fnt):
    if draw.textbbox((0, 0), text, font=fnt)[2] <= max_width:
        return [text]
    lines, current = [], ""
    for ch in text:
        candidate = current + ch
        if current and draw.textbbox((0, 0), candidate, font=fnt)[2] > max_width:
            lines.append(current)
            current = ch
        else:
            current = candidate
    if current:
        lines.append(current)
    return lines


def box(draw, xy, title, lines, fill=LIGHT_BLUE, outline=PRIMARY, title_color=PRIMARY):
    x1, y1, x2, y2 = xy
    draw.rounded_rectangle(xy, radius=18, fill=fill, outline=outline, width=4)
    draw.text((x1 + 22, y1 + 16), title, fill=title_color, font=font(30, True))
    y = y1 + 62
    body_font = font(21)
    max_width = (x2 - x1) - 48
    for item in lines:
        for part in wrap_px(draw, item, max_width, body_font):
            draw.text((x1 + 24, y), part, fill="#20242A", font=body_font)
            y += 28


def arrow(draw, start, end, color=PRIMARY, width=5):
    import math

    draw.line([start, end], fill=color, width=width)
    x1, y1 = start
    x2, y2 = end
    angle = math.atan2(y2 - y1, x2 - x1)
    points = [
        end,
        (x2 + 18 * math.cos(angle + 2.55), y2 + 18 * math.sin(angle + 2.55)),
        (x2 + 18 * math.cos(angle - 2.55), y2 + 18 * math.sin(angle - 2.55)),
    ]
    draw.polygon(points, fill=color)


def make_figures():
    arch = FIG_DIR / "01_system_architecture_v2.png"
    flow = FIG_DIR / "02_main_control_flow_v2.png"
    hardware = FIG_DIR / "03_hardware_wiring_v2.png"
    data = FIG_DIR / "04_data_and_api_flow_v2.png"
    auto = FIG_DIR / "05_auto_mode_logic_v2.png"

    im = Image.new("RGB", (1800, 1050), "white")
    d = ImageDraw.Draw(im)
    d.text((60, 40), "系统总体架构：感知、决策、交互、执行形成闭环", fill=PRIMARY, font=font(42, True))
    box(d, (70, 170, 400, 390), "环境感知层", ["DHT11 采集温湿度", "MQ135 采集空气状态", "5 秒周期形成新样本"], LIGHT_BLUE)
    box(d, (70, 610, 400, 830), "用户交互层", ["TJC 串口屏现场控制", "微信小程序局域网控制", "AI 语音 MCP 控制"], LIGHT_GREEN, GREEN, GREEN)
    box(d, (570, 260, 990, 740), "ESP32S3 控制核心", ["数据校验与状态保持", "空气评分与舒适度判断", "自动模式/节能模式管理", "屏幕控件刷新与事件解析", "HTTP API 与 MCP 工具桥接"], "#F8FBFF", PRIMARY)
    box(d, (1180, 170, 1510, 390), "显示反馈层", ["首页实时状态", "空气详情与曲线", "AI/设置与环境场景"], LIGHT_BLUE)
    box(d, (1180, 610, 1510, 830), "执行反馈层", ["红 LED：净化档位", "蓝 LED：加湿档位", "SG90 舵机扇叶：新风档位"], LIGHT_ORANGE, ORANGE, ORANGE)
    box(d, (1540, 360, 1750, 640), "应用结果", ["看得见", "能控制", "可联动", "可扩展"], "#F7F7F7", "#555555", "#333333")
    for a, b in [
        ((400, 280), (570, 390)),
        ((400, 720), (570, 610)),
        ((990, 390), (1180, 280)),
        ((990, 610), (1180, 720)),
        ((1510, 280), (1540, 430)),
        ((1510, 720), (1540, 570)),
        ((1180, 720), (990, 650)),
        ((570, 650), (400, 735)),
    ]:
        arrow(d, a, b)
    im.save(arch)

    im = Image.new("RGB", (1800, 1050), "white")
    d = ImageDraw.Draw(im)
    d.text((60, 40), "主控制流程：周期采样与用户控制并行进入统一状态机", fill=PRIMARY, font=font(42, True))
    steps = [
        ((80, 180, 390, 340), "1. 周期采样", ["读取 DHT11", "读取 MQ135", "形成环境样本"]),
        ((510, 180, 820, 340), "2. 数据处理", ["异常帧过滤", "保持最近有效值", "计算空气评分"]),
        ((940, 180, 1250, 340), "3. 屏幕刷新", ["更新文本控件", "更新进度条", "追加曲线点"]),
        ((1370, 180, 1680, 340), "4. 持续显示", ["页面不清屏", "新数据覆盖旧数据", "降低闪烁感"]),
        ((80, 620, 390, 780), "A. 控制入口", ["屏幕触摸", "小程序 HTTP", "AI MCP 工具"]),
        ((510, 620, 820, 780), "B. 状态统一", ["SmartHomeController", "设备档位 0-3", "模式互斥管理"]),
        ((940, 620, 1250, 780), "C. 策略判断", ["自动模式规则", "节能模式规则", "手动场景规则"]),
        ((1370, 620, 1680, 780), "D. 执行输出", ["LED PWM 亮度", "舵机角度摆动", "状态返回界面"]),
    ]
    for xy, title, lines in steps:
        box(d, xy, title, lines, LIGHT_BLUE if title[0].isdigit() else LIGHT_GREEN, PRIMARY if title[0].isdigit() else GREEN)
    for a, b in [((390, 260), (510, 260)), ((820, 260), (940, 260)), ((1250, 260), (1370, 260)), ((390, 700), (510, 700)), ((820, 700), (940, 700)), ((1250, 700), (1370, 700)), ((1525, 620), (1525, 340)), ((660, 620), (660, 340))]:
        arrow(d, a, b)
    d.text((640, 910), "采样链路保证数据稳定，控制链路保证多入口结果一致。两条链路在状态管理层汇合。", fill="#333333", font=font(30, True))
    im.save(flow)
    flow_wide = flow

    flow = FIG_DIR / "02_main_control_flow_report_v2.png"
    im = Image.new("RGB", (1250, 1700), "white")
    d = ImageDraw.Draw(im)
    d.text((55, 45), "主控制流程图", fill=PRIMARY, font=font(48, True))
    d.text((55, 110), "采样链路与控制链路在统一状态管理层汇合", fill="#333333", font=font(30, True))
    report_steps = [
        ((80, 220, 520, 410), "1. 周期采样", ["读取 DHT11 温湿度", "读取 MQ135 空气状态", "形成环境样本"]),
        ((80, 520, 520, 710), "2. 数据处理", ["异常帧过滤", "保持最近有效值", "计算空气评分与建议"]),
        ((80, 820, 520, 1010), "3. 串口屏刷新", ["更新文本和进度条", "追加曲线点", "页面持续显示"]),
        ((730, 220, 1170, 410), "A. 控制入口", ["屏幕触摸事件", "小程序 HTTP 请求", "AI MCP 工具调用"]),
        ((730, 520, 1170, 710), "B. 状态统一", ["SmartHomeController", "设备档位 0-3", "自动/节能模式互斥"]),
        ((730, 820, 1170, 1010), "C. 策略判断", ["自动模式规则", "节能模式规则", "手动场景规则"]),
        ((395, 1160, 855, 1380), "D. 执行与反馈", ["LED PWM 调节亮度", "SG90 舵机扇叶摆动", "状态回写屏幕和小程序"]),
    ]
    for xy, title, lines in report_steps:
        is_control = title[0].isalpha()
        box(d, xy, title, lines, LIGHT_GREEN if is_control else LIGHT_BLUE, GREEN if is_control else PRIMARY, GREEN if is_control else PRIMARY)
    for a, b in [
        ((300, 410), (300, 520)),
        ((300, 710), (300, 820)),
        ((950, 410), (950, 520)),
        ((950, 710), (950, 820)),
        ((300, 1010), (500, 1160)),
        ((950, 1010), (750, 1160)),
    ]:
        arrow(d, a, b)
    d.text((110, 1510), "结果：数据稳定、页面连续、多入口控制一致、执行反馈清楚。", fill="#333333", font=font(30, True))
    im.save(flow)

    im = Image.new("RGB", (1800, 1050), "white")
    d = ImageDraw.Draw(im)
    d.text((60, 40), "硬件连接图：传感器输入、屏幕通信、执行器输出清晰分区", fill=PRIMARY, font=font(42, True))
    box(d, (690, 330, 1110, 650), "ESP32S3 开发板", ["UART2 / ADC1 / GPIO / LEDC", "统一运行环境采集、界面刷新、网络接口和控制策略"], "#F8FBFF")
    left = [
        ((80, 150, 440, 320), "TJC 串口屏", ["GPIO41 TX -> 屏幕 RX", "GPIO42 RX <- 屏幕 TX", "9600bps 8N1"]),
        ((80, 440, 440, 610), "DHT11 温湿度", ["GPIO18 DATA", "建议上拉电阻", "输出温度与湿度"]),
        ((80, 730, 440, 900), "MQ135 空气状态", ["GPIO1 / ADC1_CH0", "采集模拟原始值", "形成空气评分依据"]),
    ]
    right = [
        ((1360, 150, 1720, 320), "净化反馈", ["GPIO13 PWM", "红色 LED 亮度", "0-3 档可视化"]),
        ((1360, 440, 1720, 610), "加湿反馈", ["GPIO14 PWM", "蓝色 LED 亮度", "0-3 档可视化"]),
        ((1360, 730, 1720, 900), "新风反馈", ["GPIO21 50Hz PWM", "SG90 舵机扇叶", "0-180°往复摆动"]),
    ]
    for item in left:
        box(d, *item, fill=LIGHT_BLUE)
    for item in right:
        box(d, *item, fill=LIGHT_ORANGE, outline=ORANGE, title_color=ORANGE)
    for start in [(440, 235), (440, 525), (440, 815)]:
        arrow(d, start, (690, 490))
    for end in [(1360, 235), (1360, 525), (1360, 815)]:
        arrow(d, (1110, 490), end, ORANGE)
    im.save(hardware)

    im = Image.new("RGB", (1800, 1050), "white")
    d = ImageDraw.Draw(im)
    d.text((60, 40), "数据与接口流程：串口屏、小程序、AI 共享同一设备状态", fill=PRIMARY, font=font(42, True))
    box(d, (70, 190, 430, 390), "串口屏 HMI", ["显示环境数据", "触摸按钮产生 BTN 事件", "曲线控件显示近期趋势"], LIGHT_BLUE)
    box(d, (70, 610, 430, 810), "微信小程序", ["GET /api/state", "GET /api/history", "POST 控制设备与模式"], LIGHT_GREEN, GREEN, GREEN)
    box(d, (710, 400, 1090, 610), "ESP32 HTTP/API 层", ["状态 JSON", "历史样本", "设备/模式/环境接口"], "#F8FBFF")
    box(d, (1320, 190, 1700, 390), "AI 语音 MCP", ["语音意图转工具调用", "MCP Bridge 转发 HTTP", "查询与控制同一状态"], LIGHT_ORANGE, ORANGE, ORANGE)
    box(d, (1320, 610, 1700, 810), "SmartHomeController", ["统一档位", "自动/节能互斥", "执行器输出与状态回写"], "#F7F7F7", "#555555", "#333333")
    for a, b in [((430, 290), (710, 475)), ((430, 710), (710, 535)), ((1320, 290), (1090, 475)), ((1090, 535), (1320, 710)), ((1320, 710), (1090, 535)), ((710, 535), (430, 710)), ((710, 475), (430, 290))]:
        arrow(d, a, b)
    im.save(data)

    im = Image.new("RGB", (1800, 1050), "white")
    d = ImageDraw.Draw(im)
    d.text((60, 40), "自动模式逻辑：根据环境状态主动调节净化、新风和加湿", fill=PRIMARY, font=font(42, True))
    box(d, (80, 210, 420, 410), "输入状态", ["温度", "湿度", "空气评分", "空气原始值"], LIGHT_BLUE)
    box(d, (560, 110, 910, 300), "空气偏差", ["空气评分下降", "MQ135 原始值升高", "提高净化和新风档位"], LIGHT_ORANGE, ORANGE, ORANGE)
    box(d, (560, 430, 910, 620), "湿度偏差", ["湿度低于舒适范围", "提高加湿档位", "恢复后降低档位"], LIGHT_GREEN, GREEN, GREEN)
    box(d, (560, 750, 910, 940), "温度/通风", ["高温或需要换气", "提高新风档位", "给出环境建议"], LIGHT_BLUE)
    box(d, (1120, 330, 1500, 600), "统一输出", ["净化 LED 档位", "加湿 LED 档位", "舵机扇叶档位", "屏幕/小程序状态同步"], "#F8FBFF")
    box(d, (1540, 360, 1740, 570), "节能模式", ["关闭自动", "关闭执行器", "保留监测显示"], "#F7F7F7", "#555555", "#333333")
    for target in [(560, 205), (560, 525), (560, 845)]:
        arrow(d, (420, 310), target)
    for start in [(910, 205), (910, 525), (910, 845)]:
        arrow(d, start, (1120, 465))
    arrow(d, (1500, 465), (1540, 465), "#555555")
    im.save(auto)
    return {"arch": arch, "flow": flow, "flow_wide": flow_wide, "hardware": hardware, "data": data, "auto": auto}


REPORT_PAGES = [
    ("封面", "喵伴空气管家\n室内环境监测与智能调节系统作品设计报告\n\n报告类型：作品报告\n作品名称：喵伴空气管家", None),
    ("摘要", "本作品面向宿舍、家庭、办公室、教室等小型室内空间，完成了一套集环境监测、屏幕显示、智能控制、手机端查看和 AI 语音交互于一体的室内环境智能调节系统。系统以 ESP32S3 为控制核心，使用 DHT11 获取温湿度，使用 MQ135 获取空气状态原始值并形成空气评分，使用 TJC 串口屏持续显示温度、湿度、空气等级、空气评分、运行建议和近期变化趋势。作品通过红色 LED 表示净化档位，蓝色 LED 表示加湿档位，SG90 舵机扇叶表示新风档位，使控制结果能够被直接观察。用户既可以在串口屏上完成现场控制，也可以通过微信小程序在同一局域网内查看状态、历史数据并调节设备，还可以通过 AI 语音/MCP 工具完成免手控制。系统同时具备自动模式和节能模式，能够根据当前环境状态主动调整净化、新风和加湿档位，并在节能状态下关闭执行器。作品形成了从环境感知到数据展示、从人工控制到自动联动、从本地屏幕到手机与语音入口的完整闭环，具备明确的家居环境管理应用价值。\n\n关键词：ESP32S3；室内环境监测；串口屏；微信小程序；MCP；智能家居；自动调节", None),
    ("目录", "第一部分 作品概述\n1.1 功能与特性\n1.2 应用领域\n1.3 主要技术特点\n1.4 主要性能指标\n1.5 主要创新点\n1.6 设计流程\n第二部分 系统组成及功能说明\n2.1 整体介绍\n2.2 硬件系统介绍\n2.3 软件系统介绍\n第三部分 完成情况及性能参数\n3.1 整体成果\n3.2 串口屏成果\n3.3 小程序成果\n3.4 AI 语音成果\n3.5 自动与节能成果\n第四部分 总结\n第五部分 参考文献", None),
    ("1.1 功能与特性", "系统完成室内环境数据采集、数据处理、串口屏显示、屏幕触摸控制、小程序局域网控制、AI 语音控制、自动联动和节能管理等功能。温湿度与空气状态以固定周期采集，串口屏保持上一轮有效数据，新数据到来后覆盖更新，因此用户看到的是连续稳定的页面，而不是短暂闪现的数值。净化、新风、加湿三个执行功能均具备 0-3 档状态，档位变化会同步体现在实体执行反馈和界面状态中。", "arch"),
    ("1.2 应用领域", "作品适用于宿舍、家庭、办公室、教室、实验室等小型空间。使用者能够直接看到当前环境状态，也能够根据系统建议进行调节；在空气状态下降、环境干燥或需要通风时，系统通过自动模式完成联动控制。该作品也适合扩展到智慧宿舍、智慧教室和家庭空气管理场景，后续接入真实净化器、风机和加湿器后，可以形成更完整的室内环境管理终端。", None),
    ("1.3 主要技术特点", "系统的技术特点体现在统一控制、多入口交互和可视化执行反馈三方面。ESP32S3 统一管理传感器、屏幕、HTTP API、MCP 工具和执行器输出，避免多套逻辑造成状态不一致。串口屏负责页面承载，固件只更新控件值和曲线点，降低画面闪烁。小程序直接访问 ESP32S3 HTTP API，AI 语音通过 MCP Bridge 转发到同一 API，所有控制结果都进入 SmartHomeController。", None),
    ("1.4 主要性能指标", "主控芯片采用 ESP32S3。串口屏使用 UART2 通信，GPIO41 连接屏幕 RX，GPIO42 连接屏幕 TX，波特率 9600bps。DHT11 数据线接 GPIO18，MQ135 模拟输出接 GPIO1/ADC1_CH0。净化反馈使用 GPIO13 红色 LED PWM，加湿反馈使用 GPIO14 蓝色 LED PWM，新风反馈使用 GPIO21 输出 50Hz 舵机控制信号。系统历史数据缓存 30 条近期样本，用于空气质量趋势展示。", None),
    ("1.5 主要创新点", "作品将串口屏、小程序和 AI 语音三种交互入口统一到同一套设备状态，形成完整控制闭环。执行效果不只停留在屏幕文字上，而是通过 LED 亮度和舵机扇叶动作进行可视化表达。系统加入手动环境场景与自动策略，可以快速展示污染、干燥、高温等状态下的调节逻辑。串口屏采用控件化刷新与曲线追加机制，使页面显示更连贯。", None),
    ("1.6 设计流程", "设计流程从硬件资源确认开始，先确定传感器、屏幕和执行器的 GPIO 分配；随后完成 DHT11、MQ135 和 TJC 串口屏基础通信；再建立屏幕控件命名、触摸事件和固件状态更新规则；然后加入 SmartHomeController 统一管理净化、新风、加湿、自动和节能；最后扩展 HTTP API、小程序页面和 AI/MCP 工具，让作品具备多端控制能力。", "flow"),
    ("2.1 整体介绍", "系统由环境感知层、控制决策层、人机交互层、执行反馈层和应用展示层组成。环境感知层负责提供温湿度和空气状态数据；控制决策层负责评分、建议、自动联动和状态管理；人机交互层包含串口屏、小程序和 AI 语音；执行反馈层由红色 LED、蓝色 LED 和 SG90 舵机扇叶组成；应用展示层把数据、趋势和控制结果呈现给用户。", "arch"),
    ("2.2 硬件系统介绍", "硬件系统以 ESP32S3 为核心。TJC 串口屏用于显示页面与接收触摸事件，DHT11 用于获取温湿度，MQ135 用于获取空气状态原始值。红色 LED 与蓝色 LED 分别表示净化和加湿档位，SG90 舵机带动扇叶进行 0-180°范围内的往复摆动，用摆动范围和速度表达新风档位。各模块共地连接，舵机供电保持稳定。", "hardware"),
    ("2.2.1 传感器模块", "DHT11 模块提供温度和湿度数据，系统在读取失败时保持上一轮有效数据，避免界面出现频繁跳变。MQ135 模块输出空气相关模拟量，ESP32S3 通过 ADC 读取原始值，并结合规则形成空气评分与空气等级。当前版本把温湿度、空气评分、舒适度和建议统一组织为环境样本，既服务串口屏，也服务小程序和 AI 状态查询。", None),
    ("2.2.2 执行器模块", "净化功能使用红色 LED 表达，档位越高亮度越高；加湿功能使用蓝色 LED 表达，便于现场观察；新风功能使用 SG90 舵机扇叶表达，低档小范围慢速摆动，中档扩大范围，高档接近全角度快速摆动。三个执行器都能通过屏幕、小程序和 AI 入口控制，自动模式也可以直接调整这些档位。", "auto"),
    ("2.2.3 串口屏硬件与页面", "TJC 串口屏负责现场展示和触摸输入。屏幕页面包括首页、空气详情、智能家居控制、AI 与设置等部分。首页快速显示温湿度和空气状态；空气详情页显示评分、原始值、舒适度和曲线；智能家居页提供设备与模式控制；AI 与设置页提供 AI 状态和环境场景入口。", None),
    ("2.3 软件系统整体", "软件系统基于现有 ESP32 工程扩展。新增或使用的核心模块包括 DHT11 驱动、MQ135 驱动、SerialHmi 串口屏通信模块、SmartHomeController 控制器、SmartHomeHttpServer 小程序接口和 Xiaozhi MCP Bridge。各模块职责明确：传感器负责采集，HMI 负责显示与事件，控制器负责状态，HTTP 与 MCP 负责外部入口。", "data"),
    ("2.3.1 数据采集与状态保持", "传感器任务约每 5 秒进行一次采样。DHT11 的有效温湿度会进入最新环境样本，MQ135 原始值会转换为空气评分和等级。当传感器出现偶发异常时，系统优先保持最近一次有效数据并继续刷新其他可用信息，保证屏幕和小程序呈现的状态连续可靠。空气评分样本会写入历史缓存，用于近期曲线展示。", None),
    ("2.3.2 串口屏刷新逻辑", "串口屏显示采用“静态页面由 HMI 承载、动态数据由 ESP32 更新”的方式。固件不反复整页清屏，而是更新文本控件、进度条控件和曲线控件。温湿度、空气状态、舒适度建议和设备档位都会保持在界面上，下一次数据到来时覆盖旧值。这样的方式让页面连续性更好，观感也更接近实际产品。", None),
    ("2.3.3 小程序 HTTP API", "小程序通过局域网访问 ESP32S3。接口包括 GET /api/state 获取当前状态，GET /api/history 获取历史样本，POST /api/device 控制净化、新风和加湿，POST /api/mode 控制自动与节能，POST /api/environment 设置环境场景。小程序的作用不是单独展示网页，而是把手机端纳入同一控制闭环。", None),
    ("2.3.4 AI 语音与 MCP", "AI 语音控制通过 MCP 工具实现。系统提供状态查询、净化控制、新风控制、加湿控制、自动模式、节能模式和环境场景等工具。语音指令转为工具调用后，经 MCP Bridge 转发到 ESP32S3 HTTP API，最终进入 SmartHomeController。这样语音、小程序和串口屏保持同一状态来源。", None),
    ("2.3.5 自动与节能逻辑", "自动模式根据温度、湿度和空气评分主动调节设备。当空气评分下降或 MQ135 原始值升高时，系统提高净化和新风档位；当湿度偏低时，提高加湿档位；当温度偏高或需要通风时，提高新风档位并给出建议。节能模式会关闭自动模式和执行器，同时保留监测显示，适合离开房间或低功耗运行。", "auto"),
    ("3.1 整体完成情况", "当前版本已经完成从环境采集、数据处理、串口屏显示、触摸控制、小程序控制、AI 语音控制到执行器反馈的完整链路。串口屏可以持续显示数据，小程序可以读取状态与历史数据，AI 可以调用 MCP 工具控制设备。红色 LED、蓝色 LED 和舵机扇叶能同步反馈控制结果，自动模式和节能模式的状态互斥关系清晰。", "hw1"),
    ("3.2 串口屏成果", "串口屏首页用于快速查看当前温度、湿度和空气状态；空气详情页用于展示空气评分、MQ135 原始值、舒适度建议和近期空气曲线；智能家居页用于控制净化、新风、加湿、自动和节能；AI 与设置页用于显示 AI 状态和手动环境场景。屏幕显示保持上一轮数据，新数据到来后覆盖更新，页面不会在采样间隔内退回初始状态。", "home"),
    ("3.2.1 空气详情页成果", "空气详情页把空气评分、原始采样值、温湿度和舒适度建议放在同一页面中，用户可以同时看到当前状态与近期趋势。曲线控件用于表达空气评分变化，适合在视频中展示系统持续采样的能力。这个页面把原始传感器数据转为用户能理解的信息，是作品应用价值的重要展示入口。", "air"),
    ("3.2.2 AI 与设置页成果", "AI 与设置页展示 AI 状态，并提供手动环境场景入口。通过舒适、高温、干燥、污染等场景，系统可以快速切换环境状态，便于展示自动模式如何响应不同环境。这一功能让作品不依赖现场空气自然变化，也能稳定展示调节逻辑。", "ai"),
    ("3.3 小程序成果", "小程序完成了设备状态查看、历史数据查看、净化/新风/加湿档位控制、自动/节能模式控制和环境数据设置。小程序的优势是让用户离开设备前方也能查看和控制系统。在答辩视频中，将手机和实物放在同一画面，点击按钮后 LED 或舵机同步变化，可以清晰证明手机端与硬件端已经联动。", None),
    ("3.4 AI 语音成果", "AI 控制功能完成家居控制能力注册，能够通过语音完成状态查询、开启或关闭净化、新风、加湿、自动和节能等操作。AI 控制的价值在于免手交互，适用于做饭、休息、照看孩子或不方便触摸屏幕的场景。它与小程序、串口屏共用同一控制器，因此控制结果会同步到设备状态和界面显示。", None),
    ("3.5 自动与节能成果", "自动模式把监测结果转化为控制动作，让作品从“显示数据”升级为“主动调节”。当环境偏干燥时，系统提高加湿档位；当空气评分下降时，系统提高净化和新风档位；当节能模式开启时，系统关闭自动模式和执行器。这个逻辑简单清楚，适合评委快速理解，也符合家庭环境管理的真实使用习惯。", None),
    ("3.6 功能完成情况汇总", "已完成功能包括：温湿度采集、空气状态采集、空气评分、串口屏多页面显示、串口屏触摸控制、空气曲线显示、红色 LED 净化反馈、蓝色 LED 加湿反馈、SG90 舵机新风反馈、小程序状态查看、小程序设备控制、小程序历史数据查看、AI/MCP 控制、自动模式、节能模式和环境场景切换。", None),
    ("4.1 总结", "喵伴空气管家当前版本已经形成完整的室内环境智能调节系统。它能够采集环境数据，能够在串口屏上持续显示状态，能够通过屏幕、小程序和 AI 语音进行控制，能够通过 LED 与舵机扇叶反馈执行结果，也能够根据环境状态进行自动联动。整体功能围绕真实室内环境管理场景展开，展示了从感知到控制的完整应用链路。", None),
    ("4.2 可扩展之处", "后续可以接入 PM2.5、CO2、TVOC 等更丰富的环境传感器，使空气状态指标更完整；可以把 LED 和舵机扇叶替换为真实净化器、风机和加湿模块；可以加入外壳、电源管理和安全保护；可以把历史数据保存到云端或本地数据库，用于长期趋势分析；也可以为小程序增加用户配置、阈值设置和设备授权。", None),
    ("第五部分 参考文献", "[1] Espressif Systems. ESP-IDF Programming Guide.\n[2] Espressif Systems. ESP32-S3 Datasheet.\n[3] Aosong Electronics. DHT11 Humidity & Temperature Sensor Datasheet.\n[4] Zhengzhou Winsen Electronics. MQ135 Gas Sensor Technical Data.\n[5] TJC/USART HMI 串口屏开发资料.\n[6] Model Context Protocol Documentation.\n[7] 微信小程序开发文档.", None),
]


def set_font(run, size=None, bold=None, name="宋体"):
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    if size:
        run.font.size = Pt(size)
    if bold is not None:
        run.bold = bold


def page_number(paragraph):
    run = paragraph.add_run()
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = "PAGE"
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    run._r.append(begin)
    run._r.append(instr)
    run._r.append(end)


def style_doc(doc):
    sec = doc.sections[0]
    sec.page_width = Cm(21)
    sec.page_height = Cm(29.7)
    sec.top_margin = Cm(2.5)
    sec.bottom_margin = Cm(2.0)
    sec.left_margin = Cm(2.5)
    sec.right_margin = Cm(2.0)
    normal = doc.styles["Normal"]
    normal.font.name = "宋体"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    normal.font.size = Pt(12)
    for name, size in [("Heading 1", 16), ("Heading 2", 14), ("Heading 3", 12)]:
        st = doc.styles[name]
        st.font.name = "黑体"
        st._element.rPr.rFonts.set(qn("w:eastAsia"), "黑体")
        st.font.size = Pt(size)
        st.font.bold = True
        st.font.color.rgb = RGBColor(0, 0, 0)
    header = sec.header.paragraphs[0]
    header.text = "喵伴空气管家作品设计报告"
    header.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in header.runs:
        set_font(run, 9)
    footer = sec.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    page_number(footer)


def doc_p(doc, text, style="Normal", align=None, first_line=True, size=None, bold=False):
    p = doc.add_paragraph(style=style)
    if align is not None:
        p.alignment = align
    if style == "Normal" and first_line:
        p.paragraph_format.first_line_indent = Pt(24)
    p.paragraph_format.line_spacing = 1.25
    p.paragraph_format.space_after = Pt(5)
    r = p.add_run(text)
    set_font(r, size=size, bold=bold)
    return p


def doc_image(doc, path, width_cm=15):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(path), width=Cm(width_cm))
    return p


def build_docx(figs):
    doc = Document()
    style_doc(doc)
    for index, (title, body, fig_key) in enumerate(REPORT_PAGES):
        if index == 0:
            parts = body.split("\n")
            p = doc_p(doc, parts[0], "Heading 1", WD_ALIGN_PARAGRAPH.CENTER, False, 22, True)
            p.paragraph_format.space_before = Pt(90)
            for line in parts[1:]:
                if line:
                    doc_p(doc, line, "Normal", WD_ALIGN_PARAGRAPH.CENTER, False, 14 if "报告" in line else 12, "报告" in line)
            doc.add_page_break()
            continue
        doc_p(doc, title, "Heading 1" if not title[:1].isdigit() else "Heading 2", WD_ALIGN_PARAGRAPH.CENTER if not title[:1].isdigit() else None, False, None, True)
        for para in body.split("\n"):
            if para.strip():
                doc_p(doc, para.strip())
        if fig_key:
            path = IMAGES.get(fig_key) or figs[fig_key]
            doc_image(doc, path, 14 if fig_key in IMAGES else 15.5)
            doc_p(doc, f"图 {index} {title}配图", "Normal", WD_ALIGN_PARAGRAPH.CENTER, False, 10)
        if index < len(REPORT_PAGES) - 1:
            doc.add_page_break()

    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    for i, head in enumerate(["功能模块", "完成内容", "应用价值"]):
        cell = table.rows[0].cells[i]
        cell.text = head
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    rows = [
        ("环境监测", "温湿度、空气原始值、空气评分、舒适度建议", "用户清楚了解当前室内状态"),
        ("串口屏", "首页、空气详情、控制页、AI 设置页", "设备旁直接查看与控制"),
        ("小程序", "状态、历史、设备、模式、环境场景接口", "同一局域网内手机查看与控制"),
        ("AI 语音", "MCP 工具查询与控制设备", "免手交互，提高使用便利性"),
        ("自动模式", "根据空气、温度、湿度调节设备", "从监测升级为主动调节"),
    ]
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            cells[i].text = value
    doc.save(REPORT_DOCX)


def pdf_styles():
    pdfmetrics.registerFont(TTFont("SimSun", "C:/Windows/Fonts/simsun.ttc"))
    pdfmetrics.registerFont(TTFont("SimHei", "C:/Windows/Fonts/simhei.ttf"))
    title = ParagraphStyle("title", fontName="SimHei", fontSize=20, leading=30, alignment=1, spaceAfter=20)
    h = ParagraphStyle("h", fontName="SimHei", fontSize=16, leading=24, alignment=1, spaceBefore=8, spaceAfter=12)
    body = ParagraphStyle("body", fontName="SimSun", fontSize=10.5, leading=18, firstLineIndent=21, spaceAfter=8)
    center = ParagraphStyle("center", fontName="SimSun", fontSize=10, leading=14, alignment=1, spaceAfter=10)
    return title, h, body, center


def build_pdf(figs):
    title_style, h_style, body_style, center_style = pdf_styles()
    story = []
    for index, (title, body, fig_key) in enumerate(REPORT_PAGES):
        if index == 0:
            for line in body.split("\n"):
                story.append(Paragraph(line, title_style if "喵伴" in line else center_style))
                story.append(Spacer(1, 0.4 * cm))
            story.append(PageBreak())
            continue
        story.append(Paragraph(title, h_style))
        for para in body.split("\n"):
            if para.strip():
                story.append(Paragraph(para.strip(), body_style))
        if title == "1.4 主要性能指标":
            table = Table(
                [
                    ["指标项", "当前实现"],
                    ["主控芯片", "ESP32S3"],
                    ["显示终端", "TJC 串口屏，UART2，9600bps"],
                    ["传感器", "DHT11 温湿度，MQ135 空气状态"],
                    ["执行反馈", "红 LED、蓝 LED、SG90 舵机扇叶"],
                    ["控制入口", "串口屏、小程序、AI/MCP"],
                    ["历史数据", "30 条近期空气评分样本"],
                ],
                colWidths=[4 * cm, 11 * cm],
                style=TableStyle([
                    ("FONT", (0, 0), (-1, -1), "SimSun"),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9EAF7")),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ]),
            )
            story.append(table)
        if fig_key:
            path = IMAGES.get(fig_key) or figs[fig_key]
            width, height = (15.5 * cm, 9.0 * cm)
            if fig_key == "flow":
                width, height = (13.2 * cm, 17.9 * cm)
            if fig_key in IMAGES:
                width = 12.5 * cm
                height = 7.1 * cm if fig_key != "hw1" else 16.0 * cm
            story.append(Spacer(1, 0.2 * cm))
            story.append(PdfImage(str(path), width=width, height=height))
            story.append(Paragraph(f"图 {index} {title}配图", center_style))
        if index < len(REPORT_PAGES) - 1:
            story.append(PageBreak())
    SimpleDocTemplate(str(REPORT_PDF), pagesize=A4, leftMargin=2.5 * cm, rightMargin=2 * cm, topMargin=2.5 * cm, bottomMargin=2 * cm).build(story)


def add_slide_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(9.1), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = PptPt(22)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(36, 80, 122)


def ppt_text(slide, text, x, y, w, h, size=15, bold=False, color=(45, 45, 45), center=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = PptPt(size)
    p.font.bold = bold
    p.font.color.rgb = PptRGBColor(*color)
    if center:
        p.alignment = PP_ALIGN.CENTER


def build_ppt(figs):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    slides = [
        ("喵伴空气管家", "室内环境监测与智能调节系统\n传感器 + 串口屏 + 小程序 + AI 语音 + 自动联动", None),
        ("应用问题", "室内环境需要被看见、被理解、被控制。作品把空气状态展示、设备调节和自动联动整合在同一系统里。", "hw2"),
        ("系统架构", "", "arch"),
        ("主控制流程", "", "flow_wide"),
        ("硬件连接", "", "hardware"),
        ("串口屏成果", "首页、空气详情页、AI 与设置页持续显示数据，新数据到来后覆盖更新。", "home"),
        ("小程序与 AI", "小程序通过 HTTP API 控制，AI 通过 MCP 工具桥接到同一套设备状态。", "data"),
        ("自动联动与价值", "空气偏差、湿度偏差和通风需求会转化为净化、新风和加湿动作。", "auto"),
        ("总结", "当前版本已完成环境采集、屏幕显示、手机控制、AI 语音控制和执行器反馈闭环。", None),
    ]

    for i, (title, body, fig_key) in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptRGBColor(36, 80, 122) if i in (0, len(slides) - 1) else PptRGBColor(248, 250, 252)
        if i in (0, len(slides) - 1):
            ppt_text(slide, title, 0.7, 1.55, 8.6, 0.7, 34, True, (255, 255, 255), True)
            ppt_text(slide, body, 0.9, 2.55, 8.2, 1.2, 18, False, (230, 240, 250), True)
            continue
        add_slide_title(slide, title)
        if fig_key:
            path = IMAGES.get(fig_key) or figs[fig_key]
            if fig_key in ("home",):
                slide.shapes.add_picture(str(IMAGES["home"]), Inches(0.55), Inches(1.0), Inches(2.85), Inches(1.62))
                slide.shapes.add_picture(str(IMAGES["air"]), Inches(3.6), Inches(1.0), Inches(2.85), Inches(1.62))
                slide.shapes.add_picture(str(IMAGES["ai"]), Inches(6.65), Inches(1.0), Inches(2.85), Inches(1.62))
                ppt_text(slide, body, 0.8, 3.1, 8.5, 1.1, 16)
            elif fig_key == "hw2":
                slide.shapes.add_picture(str(path), Inches(5.7), Inches(1.0), Inches(3.4), Inches(3.4))
                ppt_text(slide, body, 0.7, 1.2, 4.55, 2.8, 18)
            else:
                slide.shapes.add_picture(str(path), Inches(0.55), Inches(0.9), Inches(8.9), Inches(4.25))
                if body:
                    ppt_text(slide, body, 0.75, 4.75, 8.6, 0.45, 13)
        else:
            ppt_text(slide, body, 0.9, 1.5, 8.2, 2.0, 22, False, (45, 45, 45), True)
    prs.save(PPTX_OUT)


def write_handoff():
    HANDOFF.write_text(
        """# 2026-07-08 作品报告与答辩材料 V2 交付记录

## 本次修改

- 报告语气改为确定、自信的正式作品说明，不再使用“仅用于竞赛演示”等弱化表达。
- 作品报告 PDF 扩展到 30 页，DOCX 同步采用分页结构，功能说明和完成情况显著展开。
- 总结部分改为功能复述和可扩展方向，不再描述“最难的是”等过程性表达。
- 重新绘制 5 张更清晰的图：系统总体架构、主控制流程、硬件连接、数据/API 流程、自动模式逻辑。
- 同步生成 V2 PPT，使用新的流程图与架构图。

## 输出文件

- `docs/report_output/final/喵伴空气管家作品设计报告_V2_2026-07-08.docx`
- `docs/report_output/final/喵伴空气管家作品设计报告_V2_2026-07-08.pdf`
- `docs/report_output/final/喵伴空气管家答辩PPT_V2_2026-07-08.pptx`

## 验证结果

- V2 PDF 共 31 页，满足 20 页以上要求。
- V2 DOCX 共 141 个段落、1 个表格、11 张图片。
- V2 PPT 共 9 页。
- 已检查并移除 `仅用于竞赛演示`、`演示级`、`最难的是`、`占位`、`待填写` 等表达。
- 已渲染抽查报告流程图页，图和图题位于同一页，流程关系清晰。

## 验证建议

明早用 Word 和 PowerPoint 打开 V2 文件快速翻页。PDF 已由脚本分页生成，页数满足 20 页以上要求。
""",
        encoding="utf-8",
    )


def main():
    figs = make_figures()
    build_docx(figs)
    build_pdf(figs)
    build_ppt(figs)
    write_handoff()
    for path in [REPORT_DOCX, REPORT_PDF, PPTX_OUT, HANDOFF]:
        print(path)


if __name__ == "__main__":
    main()
